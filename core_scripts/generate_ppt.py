"""
PowerPoint 生成器 (S3)

功能: 生成最终 PPT，支持图片填充和 Markdown 列表
职责: 只负责 PPT 生成，不涉及 Word 解析或 LLM 调用 (SRP 原则)
输入: template.pptx, map.json, data.json
输出: 最终的 .pptx 文件
"""

import json
import os
import tempfile
from pptx import Presentation
from core_scripts.slide_duplicator import duplicate_slides_from_template


def set_paragraph_text(paragraph, line_text: str):
    """
    设置段落文本并处理 Markdown 列表级别

    :param paragraph: pptx 段落对象
    :param line_text: 文本内容 (可能包含前导空格表示级别)
    """
    # 计算缩进级别 (假设 2 个空格为一级)
    level = (len(line_text) - len(line_text.lstrip(" "))) // 2
    paragraph.level = level

    # 移除 Markdown 列表标记 "- "
    stripped_text = line_text.strip()
    if stripped_text.startswith("- "):
        paragraph.text = stripped_text[2:]
    else:
        paragraph.text = stripped_text


def handle_multiline_text(text_frame, value: str):
    """
    处理多行文本和 Markdown 列表

    :param text_frame: pptx TextFrame 对象
    :param value: 多行文本内容 (可能包含 Markdown 列表)
    """
    # 清空默认文本
    text_frame.clear()

    # 分割行
    lines = value.split("\n")

    # 设置第一行
    if lines:
        paragraph = text_frame.paragraphs[0]
        set_paragraph_text(paragraph, lines[0])

    # 添加后续行
    for line_text in lines[1:]:
        paragraph = text_frame.add_paragraph()
        set_paragraph_text(paragraph, line_text)


def generate_presentation(
    template_path: str, map_path: str, data_path: str, output_path: str
) -> str:
    """
    生成最终的 PowerPoint 演示文稿（两阶段方法）

    阶段 1: 使用 ZIP/XML 操作复制模板幻灯片
    阶段 2: 使用 python-pptx 填充内容

    :param template_path: PowerPoint 模板路径
    :param map_path: 模板映射文件 map.json 路径
    :param data_path: S2 输出的 data.json 路径
    :param output_path: 输出的 .pptx 文件路径
    :return: 输出文件路径
    """
    # 输入验证 (防御性编程)
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"模板文件不存在: {template_path}")
    if not os.path.exists(map_path):
        raise FileNotFoundError(f"映射文件不存在: {map_path}")
    if not os.path.exists(data_path):
        raise FileNotFoundError(f"数据文件不存在: {data_path}")

    # 加载映射和数据
    with open(map_path, "r", encoding="utf-8") as f:
        template_map = json.load(f)

    with open(data_path, "r", encoding="utf-8") as f:
        data_to_generate = json.load(f)

    # === 阶段 1: 收集要复制的幻灯片索引 ===
    slide_indices = []
    valid_slides = []  # 保存有效的 slide_data

    for slide_data in data_to_generate:
        layout_id = slide_data["layout_id"]

        # 检查布局是否存在 (防御性编程)
        if layout_id not in template_map["layouts"]:
            print(f"⚠️  警告: 布局 {layout_id} 在 map.json 中未定义，已跳过。")
            continue

        layout_info = template_map["layouts"][layout_id]

        # 获取模板幻灯片索引
        if "template_slide_index" in layout_info:
            slide_indices.append(layout_info["template_slide_index"])
            valid_slides.append(slide_data)
        else:
            raise ValueError(
                f"布局 {layout_id} 没有 template_slide_index 字段。"
                f"请重新运行 analyze_template.py 生成 map.json。"
            )

    # 使用 ZIP/XML 操作复制幻灯片
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        temp_pptx = tmp.name

    try:
        duplicate_slides_from_template(template_path, slide_indices, temp_pptx)

        # === 阶段 2: 填充内容 ===
        prs = Presentation(temp_pptx)

        # 遍历每一页幻灯片，填充内容
        for slide_idx, slide_data in enumerate(valid_slides):
            slide = prs.slides[slide_idx]
            layout_id = slide_data["layout_id"]
            content = slide_data["content"]

            layout_info = template_map["layouts"][layout_id]
            placeholder_mapping = layout_info["placeholders"]

        # 遍历填充占位符
        for key, value in content.items():
            # 检查键是否在映射中定义
            if key not in placeholder_mapping:
                print(f"⚠️  警告: LLM 返回的键 '{key}' 在 map.json 中未定义，已跳过。")
                continue

            ph_id = placeholder_mapping[key]["id"]

            # 获取占位符
            try:
                placeholder = slide.placeholders[ph_id]
            except KeyError:
                print(f"⚠️  警告: 占位符 ID {ph_id} 在布局中不存在，已跳过。")
                continue

            # 处理图片
            if "picture" in key.lower():
                if value and os.path.exists(value):
                    try:
                        placeholder.insert_picture(value)
                    except Exception as e:
                        print(f"⚠️  警告: 插入图片失败 ({value}): {e}")

            # 处理多行文本或列表
            elif isinstance(value, str) and (
                "\n" in value or value.strip().startswith("-")
            ):
                try:
                    handle_multiline_text(placeholder.text_frame, value)
                except Exception as e:
                    print(f"⚠️  警告: 处理多行文本失败: {e}")
                    placeholder.text = str(value)

            # 处理普通文本
            else:
                placeholder.text = str(value)

        # 确保输出目录存在
        output_dir = os.path.dirname(output_path)
        if output_dir and not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)

        # 保存 PowerPoint
        prs.save(output_path)

        print(f"✓ PowerPoint 生成完成，共 {len(valid_slides)} 页幻灯片")
        print(f"✓ 输出文件: {output_path}")

        return output_path

    finally:
        # 清理临时文件
        if os.path.exists(temp_pptx):
            os.unlink(temp_pptx)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="生成最终的 PowerPoint 演示文稿")
    parser.add_argument("--template", required=True, help="PowerPoint 模板路径")
    parser.add_argument("--map", required=True, help="模板映射文件 map.json 路径")
    parser.add_argument("--data", required=True, help="S2 输出的 data.json 路径")
    parser.add_argument("--output", required=True, help="输出的 .pptx 文件路径")

    args = parser.parse_args()

    try:
        generate_presentation(args.template, args.map, args.data, args.output)
    except Exception as e:
        print(f"✗ 错误: {e}")
        exit(1)
