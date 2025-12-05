"""
PPT 截图标注工具

功能：
1. 将 PPT 转换为预览图片
2. 在 PPT 截图上绘制编号圆圈
3. 支持不同状态的颜色标注

优先使用 LibreOffice 转换（高保真），如果失败则使用 python-pptx 渲染（简化预览）
"""

import io
import platform
import subprocess
from pathlib import Path
from typing import List, Dict, Optional, Tuple
from PIL import Image, ImageDraw, ImageFont


def get_soffice_path() -> Optional[str]:
    """获取 LibreOffice soffice 可执行文件路径（跨平台）"""
    system = platform.system()

    if system == "Darwin":  # macOS
        path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        return path if Path(path).exists() else None
    elif system == "Windows":
        paths = [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
        for path in paths:
            if Path(path).exists():
                return path
        return None
    else:  # Linux
        return "soffice"


def convert_ppt_to_pdf(pptx_path: Path, output_dir: Path) -> Optional[Path]:
    """
    使用 LibreOffice 将 PPT 转换为 PDF（跨平台）

    Args:
        pptx_path: PPT 文件路径
        output_dir: 输出目录

    Returns:
        生成的 PDF 文件路径，如果失败返回 None
    """
    soffice = get_soffice_path()
    if not soffice:
        return None

    output_dir.mkdir(parents=True, exist_ok=True)

    try:
        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(pptx_path),
            ],
            check=True,
            capture_output=True,
            text=True,
            timeout=120,
        )

        pdf_path = output_dir / (pptx_path.stem + ".pdf")
        if pdf_path.exists() and pdf_path.stat().st_size > 0:
            return pdf_path
        return None
    except Exception:
        return None


def convert_pdf_to_images(
    pdf_path: Path, output_dir: Path, dpi: int = 150
) -> List[Path]:
    """
    将 PDF 转换为图片

    Args:
        pdf_path: PDF 文件路径
        output_dir: 输出目录
        dpi: 图片分辨率

    Returns:
        生成的图片路径列表
    """
    from pdf2image import convert_from_path

    output_dir.mkdir(parents=True, exist_ok=True)
    images = convert_from_path(pdf_path, dpi=dpi)

    image_paths = []
    for i, image in enumerate(images, start=1):
        image_path = output_dir / f"page_{i}.png"
        image.save(image_path, "PNG")
        image_paths.append(image_path)

    return image_paths


# ============= python-pptx 渲染（后备方案）=============


def _get_shape_fill_color(shape) -> Optional[Tuple[int, int, int]]:
    """尝试获取形状的填充颜色"""
    try:
        fill = shape.fill
        if fill.type is not None:
            fore_color = fill.fore_color
            if fore_color.type is not None:
                rgb = fore_color.rgb
                if rgb:
                    return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return None


def _render_slide_to_image(prs, slide, width: int, height: int) -> Image.Image:
    """将幻灯片渲染为图片（简化版本）"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    scale_x = width / slide_width
    scale_y = height / slide_height
    scale = min(scale_x, scale_y)

    actual_width = int(slide_width * scale)
    actual_height = int(slide_height * scale)

    img = Image.new("RGB", (actual_width, actual_height), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # 背景色
    try:
        bg = slide.background
        if bg.fill.type is not None:
            bg_color = _get_shape_fill_color(bg)
            if bg_color:
                draw.rectangle([0, 0, actual_width, actual_height], fill=bg_color)
    except Exception:
        pass

    # 渲染形状
    for shape in slide.shapes:
        try:
            left = int(shape.left * scale)
            top = int(shape.top * scale)
            sw = int(shape.width * scale)
            sh = int(shape.height * scale)

            if sw < 5 or sh < 5:
                continue

            # 图片
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_shape = Image.open(io.BytesIO(shape.image.blob))
                    img_shape = img_shape.convert("RGB")
                    img_shape = img_shape.resize((sw, sh), Image.Resampling.LANCZOS)
                    img.paste(img_shape, (left, top))
                except Exception:
                    draw.rectangle(
                        [left, top, left + sw, top + sh],
                        fill=(200, 200, 200),
                        outline=(150, 150, 150),
                    )
                continue

            # 文本框
            if hasattr(shape, "text_frame") and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                fill_color = _get_shape_fill_color(shape)
                if fill_color:
                    draw.rectangle([left, top, left + sw, top + sh], fill=fill_color)
                if text:
                    try:
                        font_size = max(10, min(sh // 4, 24))
                        try:
                            font = ImageFont.truetype(
                                "/System/Library/Fonts/PingFang.ttc", font_size
                            )
                        except:
                            font = ImageFont.load_default()
                        display_text = text[:30] + "..." if len(text) > 30 else text
                        draw.text(
                            (left + 5, top + 5),
                            display_text,
                            fill=(50, 50, 50),
                            font=font,
                        )
                    except Exception:
                        pass
                continue

            # 其他形状
            fill_color = _get_shape_fill_color(shape)
            if fill_color:
                draw.rectangle(
                    [left, top, left + sw, top + sh],
                    fill=fill_color,
                    outline=(200, 200, 200),
                )
        except Exception:
            continue

    return img


def convert_ppt_to_images_fallback(
    pptx_path: Path, output_dir: Path, dpi: int = 150
) -> List[Path]:
    """
    使用 python-pptx 将 PPT 转换为简化预览图（后备方案）

    Args:
        pptx_path: PPT 文件路径
        output_dir: 输出目录
        dpi: 图片分辨率

    Returns:
        生成的图片路径列表
    """
    from pptx import Presentation

    output_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(pptx_path))

    width = int(10 * dpi)
    height = int(7.5 * dpi)

    image_paths = []
    for i, slide in enumerate(prs.slides, start=1):
        img = _render_slide_to_image(prs, slide, width, height)
        image_path = output_dir / f"page_{i}.png"
        img.save(image_path, "PNG")
        image_paths.append(image_path)

    return image_paths


def convert_ppt_to_images(
    pptx_path: Path, output_dir: Path, dpi: int = 150
) -> List[Path]:
    """
    将 PPT 转换为图片（自动选择最佳方案）

    优先使用 LibreOffice（高保真），失败时使用 python-pptx（简化预览）

    Args:
        pptx_path: PPT 文件路径
        output_dir: 输出目录
        dpi: 图片分辨率

    Returns:
        生成的图片路径列表
    """
    # 尝试 LibreOffice
    pdf_path = convert_ppt_to_pdf(pptx_path, output_dir.parent)
    if pdf_path and pdf_path.exists():
        try:
            return convert_pdf_to_images(pdf_path, output_dir, dpi)
        except Exception:
            pass

    # 后备：python-pptx 渲染
    return convert_ppt_to_images_fallback(pptx_path, output_dir, dpi)


def annotate_screenshot(
    image_path: Path,
    shapes_info: List[Dict],
    slide_width: int = 12192000,
    slide_height: int = 6858000,
) -> Path:
    """
    在 PPT 截图上标注元素编号

    使用幻灯片尺寸和图片尺寸的比例来计算坐标，而不是固定 DPI。
    这样可以确保标注位置准确，无论图片是如何生成的。

    Args:
        image_path: 截图路径
        shapes_info: 元素信息列表（已过滤背景元素）
        slide_width: 幻灯片宽度（EMU 单位）
        slide_height: 幻灯片高度（EMU 单位）

    Returns:
        标注后的图片路径
    """
    img = Image.open(image_path)
    draw = ImageDraw.Draw(img)

    # 获取图片实际尺寸
    img_width, img_height = img.size

    # 计算缩放比例
    scale_x = img_width / slide_width
    scale_y = img_height / slide_height

    # 加载字体
    try:
        # macOS 系统字体
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
    except:
        try:
            # Linux 系统字体
            font = ImageFont.truetype(
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 24
            )
        except:
            # 使用默认字体
            font = ImageFont.load_default()

    visible_idx = 0  # 只给可见元素分配编号
    for shape in shapes_info:
        # 跳过隐藏元素
        if shape.get("is_hidden"):
            continue

        visible_idx += 1

        # 使用比例计算坐标（而不是 DPI）
        left = shape["left"] * scale_x
        top = shape["top"] * scale_y
        width = shape["width"] * scale_x
        height = shape["height"] * scale_y

        # 编号位置：元素左上角
        x = int(left + 20)
        y = int(top + 20)

        # 确保编号在图片范围内
        x = max(20, min(x, img_width - 20))
        y = max(20, min(y, img_height - 20))

        # 确定颜色（根据是否已命名）
        is_named = shape.get("is_named", False)
        circle_color = "#007BFF" if is_named else "#FFC107"

        # 绘制圆形背景
        radius = 18
        draw.ellipse(
            [x - radius, y - radius, x + radius, y + radius],
            fill=circle_color,
            outline="#FFFFFF",
            width=2,
        )

        # 绘制编号（使用可见元素的序号）
        text = str(visible_idx)

        # 计算文本位置（居中）
        try:
            bbox = draw.textbbox((0, 0), text, font=font)
            text_width = bbox[2] - bbox[0]
            text_height = bbox[3] - bbox[1]
        except:
            # 如果 textbbox 不可用，使用估算
            text_width = len(text) * 12
            text_height = 18

        draw.text(
            (x - text_width // 2, y - text_height // 2 - 2),
            text,
            fill="#FFFFFF",
            font=font,
        )

        # 绘制元素边框（半透明虚线效果）
        border_color = "#007BFF" if is_named else "#FFC107"
        draw.rectangle(
            [left, top, left + width, top + height],
            outline=border_color,
            width=2,
        )

    # 保存标注后的图片
    annotated_path = image_path.parent / f"{image_path.stem}_annotated.png"
    img.save(annotated_path)
    return annotated_path
