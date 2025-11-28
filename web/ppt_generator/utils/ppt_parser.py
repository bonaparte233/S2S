"""
PPT 模板解析工具

功能：
1. 提取 PPT 元素坐标信息
2. 自动过滤背景元素
3. 判断元素是否已命名
"""

import re
from pathlib import Path
from typing import Dict, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def is_generic_name(name: str) -> bool:
    """判断是否为通用名称（未命名）"""
    generic_pattern = re.compile(
        r"^(图片|文本框|矩形|圆角|任意|椭圆|线条|组合|对象|"
        r"table|textbox|picture|group|rectangle|oval|line|object)\s*\d*$",
        re.IGNORECASE,
    )
    return bool(generic_pattern.match(name))


def is_background_element(shape, slide_width: int, slide_height: int) -> bool:
    """
    判断形状是否为背景/装饰元素（应该被过滤）

    返回 True 表示应该过滤掉
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        shape_type = shape.shape_type
    except:
        return True  # 无法识别的类型，过滤

    # 计算面积比例
    shape_area = shape.width * shape.height
    slide_area = slide_width * slide_height
    area_ratio = shape_area / slide_area if slide_area > 0 else 0

    # 超过 70% 幻灯片面积的通常是背景
    if area_ratio > 0.7:
        return True

    # 线条过滤
    if shape_type == MSO_SHAPE_TYPE.LINE:
        return True

    # 圆角矩形等装饰形状：如果没有文本内容则过滤
    if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        name = shape.name.lower() if shape.name else ""
        # 检查是否有文本
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            text = shape.text.strip() if shape.text else ""
            if len(text) > 0:
                return False  # 有文本，保留
        # 没有文本的自动形状，检查名称
        if "圆角矩形" in name or "矩形" in name or "椭圆" in name:
            return True

    return False


def extract_shapes_from_group(
    group_shape, slide_width: int, slide_height: int, parent_name: str = ""
) -> list:
    """
    递归提取 GROUP 中的可编辑元素

    Args:
        group_shape: GROUP 形状对象
        slide_width/slide_height: 幻灯片尺寸
        parent_name: 父级 GROUP 的名称（用于生成层级名称）

    Returns:
        list of shape info dicts
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    results = []
    group_name = group_shape.name if group_shape.name else ""

    # 构建当前路径名称
    current_path = f"{parent_name}/{group_name}" if parent_name else group_name

    for child in group_shape.shapes:
        try:
            child_type = child.shape_type
        except:
            continue

        # 递归处理嵌套的 GROUP
        if child_type == MSO_SHAPE_TYPE.GROUP:
            nested = extract_shapes_from_group(
                child, slide_width, slide_height, current_path
            )
            results.extend(nested)
            continue

        # 跳过背景元素
        if is_background_element(child, slide_width, slide_height):
            continue

        # 图片：保留
        if child_type == MSO_SHAPE_TYPE.PICTURE:
            # 跳过名称为"背景"的图片
            if child.name and "背景" in child.name:
                continue
            results.append(
                {"shape": child, "group_path": current_path, "type": "image"}
            )
            continue

        # 文本框：检查是否有内容
        if child_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text = child.text.strip() if hasattr(child, "text") and child.text else ""
            if len(text) > 0:
                results.append(
                    {
                        "shape": child,
                        "group_path": current_path,
                        "type": "text",
                        "text": text,
                    }
                )
            continue

        # 其他有文本的形状
        if hasattr(child, "has_text_frame") and child.has_text_frame:
            text = child.text.strip() if child.text else ""
            if len(text) > 0:
                results.append(
                    {
                        "shape": child,
                        "group_path": current_path,
                        "type": "text",
                        "text": text,
                    }
                )

    return results


def is_editable_content(shape, slide_width: int, slide_height: int) -> bool:
    """
    智能判断形状是否为可编辑内容（而非装饰性背景）
    注意：GROUP 类型需要单独处理，此函数不处理 GROUP
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        shape_type = shape.shape_type
    except:
        return False

    # GROUP 需要递归处理，这里返回 True 让调用方处理
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return True

    # 背景元素过滤
    if is_background_element(shape, slide_width, slide_height):
        return False

    # 图片：保留
    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True

    # 表格：保留
    if shape_type == MSO_SHAPE_TYPE.TABLE:
        return True

    # 文本框：检查是否有内容
    if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        text = shape.text.strip() if hasattr(shape, "text") and shape.text else ""
        return len(text) > 0

    # 占位符：保留
    if shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        return True

    # 图表：保留
    if shape_type == MSO_SHAPE_TYPE.CHART:
        return True

    # 其他有文本的形状
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        text = shape.text.strip() if shape.text else ""
        if len(text) > 0:
            return True

    return False


def extract_shapes_info(pptx_path: Path) -> Dict:
    """
    提取所有元素的坐标信息，自动过滤背景元素

    Returns:
        {
            'slide_width': 12192000,  # EMU 单位
            'slide_height': 6858000,
            'pages': [
                {
                    'page_num': 1,
                    'shapes': [
                        {
                            'shape_id': 0,
                            'shape_index': 3,  # 在 slide.shapes 中的实际索引
                            'name': 'TextBox 1',
                            'type': 'text',
                            'left': 914400,  # EMU
                            'top': 1828800,
                            'width': 4572000,
                            'height': 914400,
                            'text_sample': '人工智能导论',
                            'char_count': 6,
                            'is_named': False,
                            'is_hidden': False,
                            'z_order': 3  # 层级顺序（越大越在上层）
                        }
                    ]
                }
            ]
        }
    """
    prs = Presentation(str(pptx_path))
    pages = []

    # 获取幻灯片尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    for slide_num, slide in enumerate(prs.slides, 1):
        shapes_info = []
        shape_counter = 0  # 用于生成唯一的 shape_id

        for shape_index, shape in enumerate(slide.shapes):
            # 跳过母版占位符（但保留内容占位符）
            if shape.is_placeholder:
                # 检查占位符是否有实际内容
                if shape.has_text_frame:
                    text = shape.text.strip() if shape.text else ""
                    if not text:
                        continue
                else:
                    continue

            # 使用智能算法判断是否为可编辑内容
            if not is_editable_content(shape, slide_width, slide_height):
                continue

            # 处理 GROUP 类型：递归提取内部元素
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_shapes = extract_shapes_from_group(
                    shape, slide_width, slide_height
                )
                for gs in group_shapes:
                    child_shape = gs["shape"]
                    is_named = not is_generic_name(child_shape.name)

                    info = {
                        "shape_id": child_shape.shape_id,
                        "shape_index": shape_index,  # 父 GROUP 的索引
                        "name": child_shape.name,
                        "type": gs["type"],
                        "left": child_shape.left,
                        "top": child_shape.top,
                        "width": child_shape.width,
                        "height": child_shape.height,
                        "is_named": is_named,
                        "is_hidden": False,
                        "z_order": shape_counter,
                        "group_path": gs.get("group_path", ""),  # 记录 GROUP 路径
                    }

                    if gs["type"] == "text":
                        info["text_sample"] = gs.get("text", "")[:100]
                        info["char_count"] = len(gs.get("text", ""))

                    shapes_info.append(info)
                    shape_counter += 1
                continue

            # 判断元素类型
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_type = "image"
            elif shape.has_text_frame:
                shape_type = "text"
            else:
                shape_type = "other"

            # 只保留文本框和图片
            if shape_type not in ("text", "image"):
                continue

            # 判断是否已命名（非通用名称）
            is_named = not is_generic_name(shape.name)

            info = {
                "shape_id": shape.shape_id,  # 使用 shape 的真实 ID
                "shape_index": shape_index,  # 在 slide.shapes 中的索引
                "name": shape.name,
                "type": shape_type,
                "left": shape.left,  # EMU 单位
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "is_named": is_named,
                "is_hidden": False,  # 用户可以手动隐藏
                "z_order": shape_counter,  # 层级顺序
            }

            # 如果是文本框，提取示例文本
            if shape.has_text_frame:
                info["text_sample"] = shape.text[:100] if shape.text else ""
                info["char_count"] = len(shape.text) if shape.text else 0

            shapes_info.append(info)
            shape_counter += 1

        # 按 z_order 排序，确保上层元素在后面
        shapes_info.sort(key=lambda x: x["z_order"])

        pages.append({"page_num": slide_num, "shapes": shapes_info})

    return {
        "slide_width": slide_width,
        "slide_height": slide_height,
        "pages": pages,
    }


def find_shape_by_id(shapes, shape_id: int):
    """
    递归查找指定 shape_id 的形状（支持 GROUP 内嵌套）
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_id == shape_id:
            return shape
        # 递归搜索 GROUP
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            found = find_shape_by_id(shape.shapes, shape_id)
            if found:
                return found
    return None


def update_shape_name(
    pptx_path: Path, page_num: int, shape_id: int, new_name: str
) -> None:
    """
    更新 PPT 中指定元素的名称

    Args:
        pptx_path: PPT 文件路径
        page_num: 页码（从 1 开始）
        shape_id: 元素的 shape_id（支持 GROUP 内的元素）
        new_name: 新名称
    """
    prs = Presentation(str(pptx_path))

    if page_num < 1 or page_num > len(prs.slides):
        raise ValueError(f"页码 {page_num} 超出范围")

    slide = prs.slides[page_num - 1]

    # 使用 shape_id 查找元素（包括 GROUP 内的元素）
    shape = find_shape_by_id(slide.shapes, shape_id)
    if not shape:
        raise ValueError(f"找不到 shape_id={shape_id} 的元素")

    shape.name = new_name

    # 保存修改
    prs.save(str(pptx_path))
