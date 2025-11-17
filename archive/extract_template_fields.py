"""从模板 PPT 中抽取可填写的文本/图片区域，生成 JSON 骨架。"""

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

IGNORE_KEYWORDS = ("背景", "矩形", "圆角", "椭圆", "形状", "图形", "遮罩", "底色")


def _iter_shapes_with_path(shapes, parent_path=None):
    """递归遍历幻灯片（含组合），yield (shape, 路径列表)。"""
    for idx, shape in enumerate(shapes, start=1):
        name = (shape.name or "").strip() or f"元素{idx}"
        path = (*parent_path, name) if parent_path else (name,)
        yield shape, path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_with_path(shape.shapes, path)


def _is_picture(shape):
    """判断形状是否可视为图片占位符。"""
    name = shape.name or ""
    if "图片区" in name or name.startswith("图片"):
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.is_placeholder:
        try:
            return shape.placeholder_format.type == 18
        except ValueError:
            return False
    return False


def _detect_prefix(slide):
    """根据命名统计，推测当前页的前缀（如“图文页1”）。"""
    counter = {}
    for shape, path in _iter_shapes_with_path(slide.shapes):
        for segment in path:
            prefix = _extract_prefix(segment)
            if prefix:
                counter[prefix] = counter.get(prefix, 0) + 1
    if not counter:
        return None
    return max(counter.items(), key=lambda item: item[1])[0]


def _collect_fields(slide, prefix):
    """返回(树结构, 叶子映射)，仅保留可编辑文本/图片。"""
    tree = {}
    leaf_map = {}
    seen_paths = set()

    for shape, path in _iter_shapes_with_path(slide.shapes):
        label_path = _normalize_path(path, prefix)
        if not label_path:
            continue

        path_key = tuple(label_path)
        if path_key in seen_paths:
            continue
        seen_paths.add(path_key)

        if _is_picture(shape) or shape.has_text_frame:
            _insert_path(tree, label_path)
            leaf_map[path_key] = label_path[0]

    return tree, leaf_map


def _insert_path(tree, segments):
    node = tree
    for seg in segments[:-1]:
        child = node.get(seg)
        if isinstance(child, dict):
            node = child
        elif child is None:
            node[seg] = {}
            node = node[seg]
        else:
            node[seg] = {}
            node = node[seg]
    leaf = segments[-1]
    if leaf not in node or isinstance(node[leaf], dict):
        node[leaf] = ""


def _merge_content(tree, source, leaf_map):
    """根据 source 填充树的叶子节点。"""

    def dfs(node, path, src):
        if not isinstance(node, dict):
            return

        for key, value in node.items():
            new_path = path + (key,)
            if isinstance(value, dict):
                child_src = src.get(key) if isinstance(src, dict) else {}
                dfs(value, new_path, child_src or {})
            else:
                source_key = leaf_map.get(new_path)
                if (
                    isinstance(src, dict)
                    and source_key in src
                    and isinstance(src[source_key], str)
                ):
                    node[key] = src[source_key]
                else:
                    node[key] = ""

    dfs(tree, tuple(), source or {})


def _load_fill_map(fill_path):
    fill_map = {}
    data = json.loads(Path(fill_path).read_text(encoding="utf-8"))
    for page in data.get("ppt_pages", []):
        tpl = page.get("template_page_num")
        prefix = page.get("page_type")
        key = (tpl, prefix)
        fill_map[key] = page.get("content", {})
        # 备用键：按模板页码匹配
        fill_map.setdefault((tpl, None), page.get("content", {}))
    return fill_map


def _extract_prefix(name):
    """返回名称中的前缀（xxx_），仅保留包含“页”的部分。"""
    if "_" not in name:
        return None
    prefix = name.split("_", 1)[0]
    return prefix if "页" in prefix else None


def _normalize_path(path, page_prefix):
    """根据路径提取层级段，只保留与前缀相关的部分。"""
    start_idx = 0
    if page_prefix:
        start_idx = -1
        for idx, segment in enumerate(path):
            if page_prefix in segment:
                start_idx = idx
                break
        if start_idx == -1:
            return []
    trimmed = []
    for segment in path[start_idx:]:
        seg = _clean_segment(segment, page_prefix)
        if not seg:
            continue
        trimmed.append(seg)
    return trimmed


def _clean_segment(segment, page_prefix):
    seg = segment.strip()
    if not seg:
        return ""
    if page_prefix and seg.startswith(page_prefix + "_"):
        seg = seg[len(page_prefix) + 1 :]
    for keyword in IGNORE_KEYWORDS:
        if keyword in seg:
            return ""
    return seg


def _read_order(order_path):
    order_text = Path(order_path).read_text(encoding="utf-8")
    tokens = []
    for line in order_text.splitlines():
        line = line.strip()
        if not line or line.startswith("#"):
            continue
        tokens.extend(
            part.strip() for part in line.replace(",", " ").split() if part.strip()
        )
    slide_nums = []
    for tok in tokens:
        try:
            slide_nums.append(int(tok))
        except ValueError:
            raise ValueError(f"无法解析页码：{tok}")
    return slide_nums


def extract_structure(template_path, order_path, output_path, fill_path=None):
    prs = Presentation(template_path)
    if order_path:
        slide_nums = _read_order(order_path)
    else:
        slide_nums = list(range(1, len(prs.slides) + 1))
    if not slide_nums:
        raise ValueError("排序文件中没有找到有效的页码。")

    fill_map = _load_fill_map(fill_path) if fill_path else {}

    pages = []
    for idx, slide_num in enumerate(slide_nums, start=1):
        if slide_num < 1 or slide_num > len(prs.slides):
            raise ValueError(f"排序文件第{idx}个值超出范围：{slide_num}")
        slide = prs.slides[slide_num - 1]
        prefix = _detect_prefix(slide)
        fields, leaf_map = _collect_fields(slide, prefix)
        source = (
            fill_map.get((slide_num, prefix)) or fill_map.get((slide_num, None)) or {}
        )
        _merge_content(fields, source, leaf_map)
        pages.append(
            {
                "page_type": prefix if prefix else f"模板第{slide_num}页",
                "template_page_num": slide_num,
                "content": fields,
            }
        )

    output = {"ppt_pages": pages}
    Path(output_path).write_text(
        json.dumps(output, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    print(f"✅ 已写出 JSON 结构：{output_path}")


def main():
    parser = argparse.ArgumentParser(description="抽取模板页的可填写字段")
    parser.add_argument("--template", required=True, help="模板 PPTX 路径")
    parser.add_argument("--order", help="页码顺序 txt 文件（可选）")
    parser.add_argument("--output", default="template_structure.json", help="输出 JSON")
    parser.add_argument("--fill", help="参考内容 JSON（如 testFull），用于自动填充")
    args = parser.parse_args()

    extract_structure(args.template, args.order, args.output, args.fill)


if __name__ == "__main__":
    main()
