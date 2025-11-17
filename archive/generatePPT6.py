from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import os
import warnings
import shutil

# 忽略zipfile的重复名称警告
warnings.filterwarnings("ignore", category=UserWarning, module="zipfile")


class PPTGenerator:
    def __init__(self, template_path, area_names=None):
        """
        初始化PPT生成器
        :param template_path: PPT模板路径（PPTX格式）
        :param area_names: 需要解析的区域名称列表
        """
        # 1. 复制模板文件作为工作文件
        self.working_ppt_path = "./temp_working_template.pptx"
        shutil.copy2(template_path, self.working_ppt_path)

        # 2. 加载工作PPT
        self.prs = Presentation(self.working_ppt_path)

        # 3. 保存原始模板幻灯片引用和XML
        self.original_slides_info = []
        for slide in self.prs.slides:
            self.original_slides_info.append(
                {"slide": slide, "xml": slide._element.xml}
            )

        # 4. 自动解析模板区域坐标
        self.area_positions = self._auto_parse_template_areas(area_names)

        # 5. 存储新生成的幻灯片
        self.new_slides = []

    def _auto_parse_template_areas(self, target_area_names):
        """
        优化：同时解析文本框和图片区（原生占位框/普通命名形状）
        """
        area_positions = {}

        for slide_idx, slide_info in enumerate(self.original_slides_info):
            slide = slide_info["slide"]
            print(f"\n正在解析模板第{slide_idx + 1}页的形状...")

            for shape in slide.shapes:
                # 筛选条件：形状有名称（必须命名，与代码area_name一致）
                if not shape.name:
                    continue

                # 关键优化：不限制形状类型，只要命名匹配，就解析（兼容原生占位框/普通形状）
                shape_type = ""
                if shape.is_placeholder:
                    if shape.placeholder_format.type == 18:
                        shape_type = "【原生图片占位框】"
                    else:
                        shape_type = (
                            f"【占位框（类型{shape.placeholder_format.type}）】"
                        )

                elif shape.has_text_frame:
                    shape_type = "【文本框】"
                else:
                    shape_type = "【普通形状（图片区）】"  # 手动绘制的图片区

                # 若指定了目标区域名称，只保留匹配的形状
                if target_area_names and shape.name not in target_area_names:
                    continue

                # 解析坐标（统一转为Inches）
                left = shape.left.inches
                top = shape.top.inches
                width = shape.width.inches
                height = shape.height.inches

                # 使用形状名称作为键名
                # key = f"{shape.name}"
                # area_positions[key] = (Inches(left), Inches(top), Inches(width), Inches(height), slide_idx+1)
                # print(f"  ✅ 已解析：键名={key} → 形状名称={shape.name}，页码={slide_idx+1}")
                # 新增：解析文本框字体样式
                font_info = None
                if shape.has_text_frame and shape.text_frame.paragraphs:
                    # 获取第一个段落的第一个_run（实际文本片段）
                    para = shape.text_frame.paragraphs[0]
                    # 确保段落有实际文本内容（run对象）
                    if para.runs:
                        run = para.runs[0]  # 取第一个文本片段
                        font = run.font

                        # 安全获取颜色（处理未设置颜色的情况）
                        font_color = None
                        try:
                            if font.color.rgb:
                                font_color = f"RGB({font.color.rgb[0]}, {font.color.rgb[1]}, {font.color.rgb[2]})"
                            else:
                                font_color = "默认颜色"
                        except AttributeError:
                            font_color = "未设置颜色"

                        # 提取字体信息（处理None值）
                        font_info = {
                            "name": font.name if font.name else "默认字体",
                            "size": f"{font.size.pt}pt" if font.size else "默认大小",
                            "color": font_color,
                            "bold": font.bold,
                            "italic": font.italic,
                            "underline": font.underline,
                        }
                    else:
                        font_info = {"提示": "段落中没有实际文本内容"}

                # 使用形状名称作为键名，增加字体信息
                key = f"{shape.name}"

                area_positions[key] = (
                    Inches(left),
                    Inches(top),
                    Inches(width),
                    Inches(height),
                    slide_idx + 1,
                    font_info,
                )
                print(
                    f"  ✅ 已解析：键名={key} → 形状名称={shape.name}，页码={slide_idx + 1}"
                )

                if font_info:
                    if "提示" in font_info.keys():
                        print("font_info 内容：", font_info)
                    else:
                        print(f"    字体信息：")
                        print(f"      字体名称：{font_info['name']}")
                        print(f"      字体大小：{font_info['size']}")
                        print(f"      字体颜色：{font_info['color']}")
                        print(
                            f"      粗体：{font_info['bold']}，斜体：{font_info['italic']}"
                        )

        # 验证目标区域是否全部解析
        if target_area_names:
            missing_areas = [
                name for name in target_area_names if name not in area_positions
            ]
            if missing_areas:
                print(
                    f"\n 警告：以下区域未找到，请检查模板形状名称是否一致：{missing_areas}"
                )
            else:
                print(f"\n 成功解析所有{len(target_area_names)}个区域！")
        else:
            print(f"\n 共解析到{len(area_positions)}个已命名形状")

        print("已解析的所有区域名称：", list(area_positions.keys()))

        return area_positions

    def _print_slide_shapes_info(self, slide, slide_title):
        """
        打印幻灯片所有形状的详细信息
        """
        print(f"\n{'=' * 60}")
        print(f"{slide_title} - 共{len(slide.shapes)}个形状")
        print(f"{'=' * 60}")

        for i, shape in enumerate(slide.shapes):
            shape_info = []
            shape_info.append(f"形状{i + 1}:")
            shape_info.append(f"  名称: '{shape.name if shape.name else '【未命名】'}'")
            shape_info.append(f"  类型: {self._get_shape_type(shape)}")
            shape_info.append(
                f"  位置: ({shape.left.inches:.2f}, {shape.top.inches:.2f})"
            )
            shape_info.append(
                f"  尺寸: {shape.width.inches:.2f} × {shape.height.inches:.2f}"
            )

            if shape.has_text_frame:
                text = (
                    shape.text_frame.text.strip() if shape.text_frame.text else "【空】"
                )
                shape_info.append(f"  文本: '{text}'")
                shape_info.append(f"  段落数: {len(shape.text_frame.paragraphs)}")

            if shape.is_placeholder:
                shape_info.append(f"  占位符类型: {shape.placeholder_format.type}")

            print("\n".join(shape_info))
            print("-" * 40)

    def _get_shape_type(self, shape):
        """获取形状类型的描述"""
        if shape.is_placeholder:
            return "占位符"
        elif shape.has_text_frame:
            return "文本框"
        elif hasattr(shape, "shape_type"):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return "图片"
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return "自动形状"
            elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
                return "线条"
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                return "表格"
            else:
                return f"其他({shape.shape_type})"
        else:
            return "未知形状"

    def _copy_template_slide(self, template_page_num):
        """
        复制模板中的指定页到PPT末尾
        """
        template_slide_idx = template_page_num - 1
        if not (0 <= template_slide_idx < len(self.original_slides_info)):
            raise ValueError(
                f"模板第{template_page_num}页不存在（共{len(self.original_slides_info)}页）"
            )

        print(f"\n【复制模板第{template_page_num}页到新页面】")

        # 获取模板幻灯片信息
        template_info = self.original_slides_info[template_slide_idx]

        # 打印模板页的详细形状信息
        self._print_slide_shapes_info(
            template_info["slide"], f"模板第{template_page_num}页"
        )

        # 创建一个新的空白演示文稿来存储单页
        temp_prs = Presentation()

        final_slide = self._manual_copy_slide(template_info["slide"])

        # 添加到新幻灯片列表
        self.new_slides.append(final_slide)

        # 打印新页面的详细形状信息
        self._print_slide_shapes_info(
            final_slide, f"新页面（复制模板第{template_page_num}页后）"
        )

        return final_slide

    def _manual_copy_slide(self, template_slide):
        """手动复制幻灯片的所有形状"""
        # 创建一个空白幻灯片
        slide_layout = self.prs.slide_layouts[0]
        new_slide = self.prs.slides.add_slide(slide_layout)

        # 复制所有形状
        for template_shape in template_slide.shapes:
            self._copy_single_shape(template_shape, new_slide)

        return new_slide

    def _copy_single_shape(self, source_shape, target_slide):
        """复制单个形状"""
        try:
            if source_shape.has_text_frame:
                # 复制文本框
                new_shape = target_slide.shapes.add_textbox(
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height,
                )

                # 复制文本内容
                if source_shape.text_frame.text:
                    new_shape.text_frame.text = source_shape.text_frame.text

                # 复制名称
                if source_shape.name:
                    new_shape.name = source_shape.name

                print(
                    f"  ✅ 复制文本框: {source_shape.name if source_shape.name else '未命名'}"
                )

            elif hasattr(source_shape, "image") and source_shape.image:
                # 复制图片
                try:
                    img_bytes = source_shape.image.blob
                    temp_img_path = f"./temp_img_{id(source_shape)}.png"
                    with open(temp_img_path, "wb") as f:
                        f.write(img_bytes)

                    new_picture = target_slide.shapes.add_picture(
                        temp_img_path,
                        source_shape.left,
                        source_shape.top,
                        source_shape.width,
                        source_shape.height,
                    )
                    if source_shape.name:
                        new_picture.name = source_shape.name
                    os.remove(temp_img_path)
                    print(
                        f"  ✅ 复制图片: {source_shape.name if source_shape.name else '未命名'}"
                    )
                except Exception as e:
                    print(f"  ❌ 复制图片失败: {e}")

            else:
                # 对于其他形状，创建文本框占位
                new_shape = target_slide.shapes.add_textbox(
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height,
                )
                if source_shape.name:
                    new_shape.name = source_shape.name
                print(
                    f"  ✅ 复制其他形状: {source_shape.name if source_shape.name else '未命名'}"
                )

        except Exception as e:
            print(f"  ❌ 复制形状失败: {e}")

    def _get_installed_fonts(self):
        """获取系统已安装字体列表"""
        installed_fonts = set()
        # Windows系统字体注册表路径
        reg_path = r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts"

        try:
            with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, reg_path) as key:
                i = 0
                while True:
                    try:
                        font_name, _, _ = winreg.EnumValue(key, i)
                        # 提取字体名称（去除版本等信息）
                        clean_name = re.sub(
                            r" \(TrueType\)|\.ttf$", "", font_name, flags=re.IGNORECASE
                        )
                        installed_fonts.add(clean_name.lower())
                        i += 1
                    except OSError:
                        break
        except Exception as e:
            print(f"获取已安装字体失败: {e}")

        return installed_fonts

    def _install_font(self, font_path):
        """安装字体到系统"""
        try:
            # Windows系统字体目录
            font_dir = os.path.join(os.environ["SystemRoot"], "Fonts")
            font_filename = os.path.basename(font_path)
            dest_path = os.path.join(font_dir, font_filename)

            # 复制字体文件到系统字体目录
            shutil.copy2(font_path, dest_path)

            # 注册字体到注册表
            font_name = os.path.splitext(font_filename)[0]
            reg_path = r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts"
            with winreg.OpenKey(
                winreg.HKEY_LOCAL_MACHINE, reg_path, 0, winreg.KEY_SET_VALUE
            ) as key:
                winreg.SetValueEx(key, font_name, 0, winreg.REG_SZ, font_filename)

            # 通知系统字体更新
            ctypes.windll.gdi32.AddFontResourceW(dest_path)
            ctypes.windll.user32.SendMessageW(
                wintypes.HWND_BROADCAST,
                0x001D,  # WM_FONTCHANGE
                0,
                0,
            )
            print(f"✅ 字体安装成功: {font_name}")
            return True
        except Exception as e:
            print(f"❌ 字体安装失败: {e}")
            return False

    def _download_font(self, font_name):
        """从网络下载字体（这里使用示例API，实际需替换为可靠源）"""
        try:
            # 注意：实际应用中需要使用可靠的字体下载源
            # 以下仅为示例，可能无法直接使用
            print(f"正在尝试下载字体: {font_name}")
            url = f"https://example.com/fonts/{font_name.replace(' ', '+')}.ttf"  # 示例URL

            with tempfile.NamedTemporaryFile(suffix=".ttf", delete=False) as tmp_file:
                response = requests.get(url, timeout=10)
                response.raise_for_status()
                tmp_file.write(response.content)
                tmp_file_path = tmp_file.name

            # 验证字体文件
            try:
                TTFont(tmp_file_path)
                print(f"✅ 字体文件验证成功")
                return tmp_file_path
            except:
                print(f"❌ 下载的字体文件无效")
                os.remove(tmp_file_path)
                return None

        except Exception as e:
            print(f"❌ 字体下载失败: {e}")
            return None

    def _ensure_font_available(self, font_name):
        """确保字体可用，不存在则下载安装"""
        if not font_name:
            return None

        installed_fonts = self._get_installed_fonts()
        font_name_lower = font_name.lower()

        # 检查字体是否已安装
        if any(font_name_lower in font.lower() for font in installed_fonts):
            print(f"字体已安装: {font_name}")
            return font_name

        # 尝试下载并安装字体
        print(f"未找到字体: {font_name}")
        font_path = self._download_font(font_name)
        if font_path and self._install_font(font_path):
            os.remove(font_path)  # 清理临时文件
            return font_name
        return None

    # def add_text_to_area(self, slide, template_page_num, area_name, text, font_size=None, font_color=None):
    # """向指定模板页的文本框添加文字，支持字体样式继承"""
    # # 使用形状名称作为键名
    # area_key = f"{area_name}"
    # print(f"\n正在查找区域：{area_key}")

    # if area_key not in self.area_positions:
    # print(f"已解析的所有区域键：{list(self.area_positions.keys())}")
    # raise ValueError(f"❌ 未找到「{area_name}」区域，请检查模板形状名称")

    # # 解析区域信息，包含字体信息
    # target_left, target_top, target_width, target_height, _, font_info = self.area_positions[area_key]
    # found_textbox = None

    # # 查找当前幻灯片中匹配的文本框（按名称+位置）
    # for shape in slide.shapes:
    # if (shape.name == area_name and
    # shape.has_text_frame and
    # abs(shape.left.inches - target_left.inches) < 0.01 and
    # abs(shape.top.inches - target_top.inches) < 0.01):
    # found_textbox = shape
    # break
    # # if found_textbox:
    # # text_frame = found_textbox.text_frame
    # # text_frame.word_wrap = True

    # # # 关键修复：不清空文本框，而是清空第一个段落的内容（保留段落结构和格式）
    # # if text_frame.paragraphs:
    # # # 清空第一个段落的所有文本片段（run）
    # # para = text_frame.paragraphs[0]
    # # para.clear()  # 只清空内容，保留段落格式
    # # else:
    # # # 如果没有段落，才新建一个
    # # para = text_frame.add_paragraph()

    # # # 直接在保留格式的段落中添加文本
    # # para.text = text

    # # # 应用字体格式（此时para.font会继承原始段落的格式）
    # # font = para.font  # 这里的font会基于原始段落格式

    # # # 处理字体（使用模板字体信息覆盖）
    # # if font_info and font_info['name'] and "默认" not in font_info['name']:
    # # font.name = font_info['name']

    # # # 处理字号（优先使用传入的font_size，否则用模板信息）
    # # if font_size:
    # # font.size = Pt(font_size)
    # # elif font_info and font_info['size'] and "默认" not in font_info['size']:
    # # try:
    # # size_str = font_info['size'].split('pt')[0].strip()
    # # size_value = float(size_str)
    # # font.size = Pt(size_value)
    # # print(f"应用字体大小: {size_value}pt")  # 打印实际应用的磅值
    # # except (ValueError, TypeError):
    # # print(f"⚠️ 无法解析字体大小: {font_info['size']}，使用默认大小")

    # # # 处理字体颜色（补充颜色应用逻辑，原代码缺失）
    # # if font_color:
    # # # 假设font_color是"RGB(r,g,b)"格式，解析并应用
    # # try:
    # # r, g, b = map(int, font_color.strip('RGB()').split(','))
    # # font.color.rgb = RGBColor(r, g, b)
    # # except:
    # # print(f"⚠️ 无法解析字体颜色: {font_color}")
    # # elif font_info and font_info['color'] and "默认" not in font_info['color']:
    # # try:
    # # r, g, b = map(int, font_info['color'].strip('RGB()').split(','))
    # # font.color.rgb = RGBColor(r, g, b)
    # # except:
    # # print(f"⚠️ 无法解析模板字体颜色: {font_info['color']}")

    # # # 处理粗体和斜体（补充原代码缺失的逻辑）
    # # if font_info:
    # # if font_info['bold'] is not None:
    # # font.bold = font_info['bold']
    # # if font_info['italic'] is not None:
    # # font.italic = font_info['italic']
    # # print(f"✅ 文本填充成功：{area_name}")
    # if found_textbox:
    # text_frame = found_textbox.text_frame
    # text_frame.word_wrap = True

    # # # 清空现有文本
    # # text_frame.clear()

    # if text_frame.paragraphs:
    # # 清空第一个段落的所有文本片段（run）
    # paragraph = text_frame.paragraphs[0]
    # paragraph.clear()  # 只清空内容，保留段落格式
    # else:
    # # 如果没有段落，才新建一个
    # paragraph = text_frame.add_paragraph()

    # # 添加新文本
    # #paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    # paragraph.text = text

    # # 应用字体格式 - 优先使用模板中的样式
    # font = paragraph.font

    # print(font_info)
    # # if font_info:
    # # print(f"    字体信息：")
    # # print(f"      字体名称：{font_info['name']}")
    # # print(f"      字体大小：{font_info['size']}")
    # # print(f"      字体颜色：{font_info['color']}")
    # # print(f"      粗体：{font_info['bold']}，斜体：{font_info['italic']}")

    # # 处理字体
    # if font_info and font_info['name'] and "默认" not in font_info['name']:
    # font.name = font_info['name']

    # # 处理字号
    # if font_size:
    # font.size = Pt(font_size)
    # # elif font_info and font_info['size']:
    # # font.size = Pt(font_info['size'])
    # elif font_info and font_info['size'] and "默认" not in font_info['size']:
    # try:
    # # 提取第一个有效的数值（处理重复pt的情况）
    # size_str = font_info['size'].split('pt')[0].strip()
    # size_value = float(size_str)
    # font.size = Pt(size_value)

    # print(f"字体大小设置为: {size_value}pt")  # 这里改为打印原始数值
    # except (ValueError, TypeError):
    # print(f"⚠️ 无法解析字体大小: {font_info['size']}，使用默认大小")

    # # # 处理颜色（安全处理None值）
    # # if font_color:
    # # font.color.rgb = font_color
    # # elif font_info and font_info['color']:
    # # font.color.rgb = font_info['color']
    # # 处理颜色（修正后）
    # if font_color:
    # # 如果传入的是RGBColor对象直接使用
    # if isinstance(font_color, RGBColor):
    # font.color.rgb = font_color
    # else:
    # # 假设传入的是"RGB(r,g,b)"格式字符串
    # try:
    # r, g, b = map(int, font_color.replace("RGB(", "").replace(")", "").split(","))
    # font.color.rgb = RGBColor(r, g, b)
    # except:
    # print(f"⚠️ 颜色格式错误: {font_color}，使用默认颜色")
    # elif font_info and font_info['color'] and "默认" not in font_info['color']:
    # # 解析模板中的颜色字符串（如"RGB(255, 0, 0)"）
    # try:
    # # 提取RGB数值
    # color_str = font_info['color']
    # if color_str.startswith("RGB(") and color_str.endswith(")"):
    # rgb_values = color_str[4:-1].split(",")
    # r = int(rgb_values[0].strip())
    # g = int(rgb_values[1].strip())
    # b = int(rgb_values[2].strip())
    # font.color.rgb = RGBColor(r, g, b)
    # except (ValueError, IndexError):
    # print(f"⚠️ 无法解析颜色: {font_info['color']}，使用默认颜色")

    # # 处理其他样式
    # if font_info:
    # font.bold = font_info['bold']
    # font.italic = font_info['italic']
    # font.underline = font_info['underline']

    # print(f"✅ 文本填充成功：{area_name}")

    # else:
    # print(f"⚠️ 未找到匹配的文本框：{area_name}")

    def add_text_to_area(
        self, slide, template_page_num, area_name, text, font_size=None, font_color=None
    ):
        """向指定模板页的文本框添加文字"""
        # 使用形状名称作为键名
        area_key = f"{area_name}"
        print(f"\n正在查找区域：{area_key}")

        if area_key not in self.area_positions:
            print(f"已解析的所有区域键：{list(self.area_positions.keys())}")
            raise ValueError(f"❌ 未找到「{area_name}」区域，请检查模板形状名称")

        target_left, target_top, target_width, target_height, _, font_info = (
            self.area_positions[area_key]
        )
        # target_left, target_top, target_width, target_height, _ = self.area_positions[area_key]
        found_textbox = None

        # 查找当前幻灯片中匹配的文本框（按名称+位置）
        for shape in slide.shapes:
            if (
                shape.name == area_name
                and shape.has_text_frame
                and abs(shape.left.inches - target_left.inches) < 0.01
                and abs(shape.top.inches - target_top.inches) < 0.01
            ):
                found_textbox = shape
                break

        if found_textbox:
            text_frame = found_textbox.text_frame
            text_frame.word_wrap = True

            # 清空现有文本
            text_frame.clear()

            # 添加新文本
            paragraph = (
                text_frame.paragraphs[0]
                if text_frame.paragraphs
                else text_frame.add_paragraph()
            )
            paragraph.text = text

            # 应用字体格式
            font = paragraph.font
            if font_size:
                font.size = Pt(font_size)
            if font_color:
                font.color.rgb = font_color

            print(f"✅ 文本填充成功：{area_name}")
        else:
            print(f"⚠️ 未找到匹配的文本框：{area_name}")

    def add_image_to_area(self, slide, template_page_num, area_name, img_path):
        """向指定模板页的图片区插入图片"""
        area_key = f"{area_name}"
        if area_key not in self.area_positions:
            raise ValueError(f"❌ 未找到「{area_name}」区域，请检查模板形状名称")

        target_left, target_top, target_width, target_height, _, font_info = (
            self.area_positions[area_key]
        )
        found_shape = None

        # 查找当前幻灯片中匹配的图片区（按名称）
        for shape in slide.shapes:
            if shape.name == area_name:
                found_shape = shape
                break

        # 校验图片路径
        if not os.path.exists(img_path):
            print(f"⚠️  警告：图片文件不存在 → 路径：{img_path}")
            return

        # 插入图片
        if found_shape:
            # 删除原有形状
            sp = found_shape._element
            sp.getparent().remove(sp)

            # 插入新图片
            slide.shapes.add_picture(
                img_path,
                left=target_left,
                top=target_top,
                width=target_width,
                height=target_height,
            )
            print(f"✅ 图片插入成功：{area_name} → {img_path}")
        else:
            print(f"⚠️ 未找到匹配的图片区域：{area_name}")

    def save_final_ppt(self, output_path):
        """保存最终PPT（只包含新生成的页面）"""
        # 创建一个全新的演示文稿
        final_prs = Presentation()

        # 复制所有新幻灯片到最终演示文稿
        for new_slide in self.new_slides:
            # 手动复制
            self._manual_copy_to_final(new_slide, final_prs)

        # 保存最终PPT
        final_prs.save(output_path)

        # 清理临时文件
        if os.path.exists(self.working_ppt_path):
            os.remove(self.working_ppt_path)

        print(f"\n🎉 PPT生成完成！路径：{output_path}")
        print(f"共生成 {len(self.new_slides)} 页PPT")

    def _manual_copy_to_final(self, source_slide, final_prs):
        """手动复制幻灯片到最终PPT"""
        slide_layout = final_prs.slide_layouts[0]
        new_slide = final_prs.slides.add_slide(slide_layout)

        # 复制所有形状
        for shape in source_slide.shapes:
            self._copy_single_shape(shape, new_slide)


def generate_ppt(template_path, script_json_path, output_path):
    # 定义所有需要解析的区域名称（与模板形状名称一致）
    required_area_names = [
        "封面页_学院名称区",
        "封面页_总课程名称区",
        "目录页_总课程名称区",
        "目录页_目录标题区",
        "目录页_目录条目区1",
        "目录页_目录条目区2",
        "目录页_目录条目区3",
        "目录页_目录条目区4",
        "主讲人页_总课程名称区",
        "主讲人页_主讲教师区",
        "过渡页_总课程名称区",
        "过渡页_文字内容区",
        "图文页1_总课程名称区",
        "图文页1_一级标题区",
        "图文页1_二级标题区",
        "图文页1_三级标题区",
        "图文页1_文字标题区",
        "图文页1_文字内容区1",
        "图文页1_文字内容区2",
        "图文页1_文字内容区3",
        "图文页1_图片区",
        "图文页2_总课程名称区",
        "图文页2_三级标题区",
        "图文页2_文字内容区1",
        "图文页2_文字内容区2",
        "图文页2_文字内容区3",
        "图文页3_总课程名称区",
        "图文页3_一级标题区",
        "图文页3_二级标题区",
        "图文页3_三级标题区",
        "图文页3_标签区1",
        "图文页3_标签区2",
        "图文页3_标签区3",
        "图文页3_标签区4",
        "图文页3_文字内容区1",
        "图文页3_文字内容区2",
        "图文页3_文字内容区3",
        "图文页3_文字内容区4",
        "图文页4_总课程名称区",
        "图文页4_一级标题区",
        "图文页4_二级标题区",
        "图文页4_三级标题区",
        "图文页4_文字内容区",
        "图文页4_图片区",
        "图文页5_总课程名称区",
        "图文页5_二级标题区",
        "图文页5_三级标题区",
        "图文页5_文字内容区1",
        "图文页5_文字内容区2",
        "图文页5_文字内容区3",
        "图文页5_图片区1",
        "图文页5_图片区2",
        "图文页5_图片区3",
        "文字页1_总课程名称区",
        "文字页1_一级标题区",
        "文字页1_二级标题区",
        "文字页1_三级标题区",
        "文字页1_文字标题区1",
        "文字页1_文字标题区2",
        "文字页1_文字内容区1-1",
        "文字页1_文字内容区1-2",
        "文字页1_文字内容区1-3",
        "文字页1_文字内容区2-1",
        "文字页1_文字内容区2-2",
        "文字页1_文字内容区2-3",
        "文字页2_总课程名称区",
        "文字页2_一级标题区",
        "文字页2_二级标题区",
        "文字页2_三级标题区",
        "文字页2_文字标题区1",
        "文字页2_文字标题区2",
        "文字页2_文字内容区1-1",
        "文字页2_文字内容区1-2",
        "文字页2_文字内容区1-3",
        "文字页2_文字内容区2-1",
        "文字页2_文字内容区2-2",
        "文字页2_文字内容区2-3",
        "文字页3_总课程名称区",
        "文字页3_一级标题区",
        "文字页3_二级标题区",
        "文字页3_三级标题区",
        "文字页3_文字标题区1",
        "文字页3_文字标题区2",
        "文字页3_文字标题区3",
        "文字页3_文字内容区1",
        "文字页3_文字内容区2",
        "文字页3_文字内容区3",
    ]

    # 初始化PPT生成器
    ppt_gen = PPTGenerator(template_path, area_names=required_area_names)

    # 读取JSON讲稿数据
    with open(script_json_path, "r", encoding="utf-8") as f:
        script_data = json.load(f)

    # 循环生成每页PPT
    for page_idx, page_data in enumerate(script_data["ppt_pages"]):
        page_type = page_data["page_type"]
        template_page_num = page_data["template_page_num"]
        content = page_data["content"]

        print(
            f"\n===== 正在生成第{page_idx + 1}页（类型：{page_type}，复用模板第{template_page_num}页）====="
        )

        # 复制模板中指定页码的幻灯片
        slide = ppt_gen._copy_template_slide(template_page_num)

        # 按页面类型填充内容
        if page_type == "封面页":
            # ppt_gen.add_text_to_area(
            # slide, template_page_num, "封面页_学院名称区", content["学院名称"],
            # 22, RGBColor(0, 82, 154)
            # )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "封面页_学院名称区",
                content["学院名称"],
                None,
                None,
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "封面页_总课程名称区",
                content["总课程名称区"],
                28,
                RGBColor(0, 82, 154),
            )

        elif page_type == "目录页":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "目录页_总课程名称区",
                content["总课程名称区"],
                20,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "目录页_目录标题区",
                content["目录标题区"],
                26,
                RGBColor(0, 82, 154),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "目录页_目录条目区1",
                content["目录条目区1"],
                26,
                RGBColor(0, 82, 154),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "目录页_目录条目区2",
                content["目录条目区2"],
                26,
                RGBColor(0, 82, 154),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "目录页_目录条目区3",
                content["目录条目区3"],
                26,
                RGBColor(0, 82, 154),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "目录页_目录条目区4",
                content["目录条目区4"],
                26,
                RGBColor(0, 82, 154),
            )

        elif page_type == "主讲人页":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "主讲人页_总课程名称区",
                content["总课程名称区"],
                20,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "主讲人页_主讲教师区",
                content["主讲教师区"],
                24,
                RGBColor(0, 82, 154),
            )

        elif page_type == "过渡页":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "过渡页_总课程名称区",
                content["总课程名称区"],
                20,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "过渡页_文字内容区",
                content["文字内容区"],
                22,
                RGBColor(51, 51, 51),
            )

        elif page_type == "图文内容页样式1":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页1_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页1_一级标题区",
                content["一级标题区"],
                24,
                RGBColor(0, 82, 154),
            )
            if "二级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页1_二级标题区",
                    content["二级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页1_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页1_文字标题区",
                content["文字标题区"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页1_文字内容区1",
                content["文字内容区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页1_文字内容区2",
                content["文字内容区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页1_文字内容区3",
                content["文字内容区3"],
                20,
                RGBColor(51, 51, 51),
            )
            if "图片区" in content and content["图片区"]:
                try:
                    ppt_gen.add_image_to_area(
                        slide, template_page_num, "图文页1_图片区", content["图片区"]
                    )
                except Exception as e:
                    print(f"⚠️  图片插入失败：{str(e)}")

        elif page_type == "图文内容页样式2":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页2_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )

            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页2_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )

            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页2_文字内容区1",
                content["文字内容区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页2_文字内容区2",
                content["文字内容区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页2_文字内容区3",
                content["文字内容区3"],
                20,
                RGBColor(51, 51, 51),
            )
            if "图片区" in content and content["图片区"]:
                try:
                    ppt_gen.add_image_to_area(
                        slide, template_page_num, "图文页2_图片区", content["图片区"]
                    )
                except Exception as e:
                    print(f"⚠️  图片插入失败：{str(e)}")

        elif page_type == "图文内容页样式3":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_一级标题区",
                content["一级标题区"],
                24,
                RGBColor(0, 82, 154),
            )
            if "二级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页3_二级标题区",
                    content["二级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页3_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_标签区1",
                content["标签区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_标签区2",
                content["标签区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_标签区3",
                content["标签区3"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_文字内容区1",
                content["文字内容区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_文字内容区2",
                content["文字内容区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页3_文字内容区3",
                content["文字内容区3"],
                20,
                RGBColor(51, 51, 51),
            )
            if "图片区" in content and content["图片区"]:
                try:
                    ppt_gen.add_image_to_area(
                        slide, template_page_num, "图文页3_图片区", content["图片区"]
                    )
                except Exception as e:
                    print(f"⚠️  图片插入失败：{str(e)}")

        elif page_type == "图文内容页样式4":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页4_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页4_一级标题区",
                content["一级标题区"],
                24,
                RGBColor(0, 82, 154),
            )
            if "二级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页4_二级标题区",
                    content["二级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页4_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )

            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页4_文字内容区",
                content["文字内容区"],
                20,
                RGBColor(51, 51, 51),
            )

            if "图片区" in content and content["图片区"]:
                try:
                    ppt_gen.add_image_to_area(
                        slide, template_page_num, "图文页4_图片区", content["图片区"]
                    )
                except Exception as e:
                    print(f"⚠️  图片插入失败：{str(e)}")

        elif page_type == "图文内容页样式5":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页5_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )

            if "二级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页5_二级标题区",
                    content["二级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "图文页5_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )

            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页5_文字内容区1",
                content["文字内容区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页5_文字内容区2",
                content["文字内容区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "图文页5_文字内容区3",
                content["文字内容区3"],
                20,
                RGBColor(51, 51, 51),
            )
            if "图片区" in content and content["图片区1"]:
                try:
                    ppt_gen.add_image_to_area(
                        slide, template_page_num, "图文页5_图片区1", content["图片区1"]
                    )
                except Exception as e:
                    print(f"⚠️  图片插入失败：{str(e)}")
            if "图片区" in content and content["图片区2"]:
                try:
                    ppt_gen.add_image_to_area(
                        slide, template_page_num, "图文页5_图片区2", content["图片区2"]
                    )
                except Exception as e:
                    print(f"⚠️  图片插入失败：{str(e)}")
            if "图片区" in content and content["图片区3"]:
                try:
                    ppt_gen.add_image_to_area(
                        slide, template_page_num, "图文页5_图片区3", content["图片区3"]
                    )
                except Exception as e:
                    print(f"⚠️  图片插入失败：{str(e)}")
        elif page_type == "文字内容页样式1":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_一级标题区",
                content["一级标题区"],
                24,
                RGBColor(0, 82, 154),
            )
            if "二级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "文字页1_二级标题区",
                    content["二级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "文字页1_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字标题区1",
                content["文字标题区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字标题区2",
                content["文字标题区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字内容区1-1",
                content["文字内容区1-1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字内容区1-2",
                content["文字内容区1-2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字内容区1-3",
                content["文字内容区1-3"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字内容区2-1",
                content["文字内容区2-1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字内容区2-2",
                content["文字内容区2-2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页1_文字内容区2-3",
                content["文字内容区2-3"],
                20,
                RGBColor(51, 51, 51),
            )
        elif page_type == "文字内容页样式2":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_一级标题区",
                content["一级标题区"],
                24,
                RGBColor(0, 82, 154),
            )
            if "二级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "文字页2_二级标题区",
                    content["二级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "文字页2_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字标题区1",
                content["文字标题区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字标题区2",
                content["文字标题区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字内容区1-1",
                content["文字内容区1-1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字内容区1-2",
                content["文字内容区1-2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字内容区1-3",
                content["文字内容区1-3"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字内容区2-1",
                content["文字内容区2-1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字内容区2-2",
                content["文字内容区2-2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页2_文字内容区2-3",
                content["文字内容区2-3"],
                20,
                RGBColor(51, 51, 51),
            )
        elif page_type == "文字内容页样式3":
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_总课程名称区",
                content["总课程名称区"],
                18,
                RGBColor(102, 102, 102),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_一级标题区",
                content["一级标题区"],
                24,
                RGBColor(0, 82, 154),
            )
            if "二级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "文字页3_二级标题区",
                    content["二级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            if "三级标题区" in content:
                ppt_gen.add_text_to_area(
                    slide,
                    template_page_num,
                    "文字页3_三级标题区",
                    content["三级标题区"],
                    22,
                    RGBColor(30, 92, 164),
                )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_文字标题区1",
                content["文字标题区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_文字标题区2",
                content["文字标题区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_文字标题区3",
                content["文字标题区3"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_文字内容区1",
                content["文字内容区1"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_文字内容区2",
                content["文字内容区2"],
                20,
                RGBColor(51, 51, 51),
            )
            ppt_gen.add_text_to_area(
                slide,
                template_page_num,
                "文字页3_文字内容区3",
                content["文字内容区3"],
                20,
                RGBColor(51, 51, 51),
            )
    # 保存最终PPT
    ppt_gen.save_final_ppt(output_path)


if __name__ == "__main__":
    # 配置文件路径
    TEMPLATE_PATH = "./template.pptx"
    # SCRIPT_JSON_PATH = "./test.json"
    SCRIPT_JSON_PATH = "./testFull.json"
    OUTPUT_PATH = "./newppt_final.pptx"

    # 执行生成
    generate_ppt(TEMPLATE_PATH, SCRIPT_JSON_PATH, OUTPUT_PATH)
