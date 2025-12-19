"""根据 JSON 描述复制模板、填充文本/图片并生成最终 PPT。"""

import argparse
import json
import posixpath
import re
import secrets
import shutil
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Dict, Optional, Set
from xml.etree import ElementTree as ET

from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


SUFFIXES = ("区", "框", "栏")
PLACEHOLDER_KEYWORDS = ("文字内容", "字幕", "标题名称", "内容内容")
SUBTITLE_TEXTS = ("字幕18pt，白色字体深色描边，悬浮阴影。确保在任何底色上都能明确显示",)
IGNORE_KEYWORDS = ("背景", "矩形", "圆角", "椭圆", "形状", "图形", "遮罩", "底色")
EXPANDABLE_KEYWORDS = ("标题", "课题", "栏目")
EMU_PER_PT = 12700
H_PADDING = 20000  # 约 1.5 毫米
MANUAL_NAME_MAP = {
    "目录内容区1": ["文本框 9"],
    "目录内容区2": ["文本框 14"],
    "目录内容区3": ["文本框 17"],
    "目录内容区4": ["文本框 20"],
    # 常见字幕框
    "字幕": ["文本框 10", "文本框 32", "文本框 36", "文本框 58", "文本框 121"],
}
NSMAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}
SLIDE_RE = re.compile(r"ppt/slides/slide(\d+)\.xml")
TAG_RE = re.compile(r"ppt/tags/tag(\d+)\.xml")

# XML 命名空间常量（用于 build_from_json）
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
OD_REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"

ET.register_namespace("", P_NS)
ET.register_namespace("r", OD_REL_NS)


def _clean_rels_namespace(xml_bytes: bytes) -> bytes:
    """清理 rels XML 中的 ns0: 前缀，PowerPoint 对此敏感。

    ElementTree 在输出时会给未注册的命名空间添加 ns0: 前缀，
    但 PowerPoint 期望 rels 文件使用无前缀的默认命名空间。
    """
    xml_str = xml_bytes.decode("utf-8")
    # 替换 ns0: 前缀和 xmlns:ns0 声明
    xml_str = xml_str.replace("ns0:", "").replace(":ns0", "")
    # 处理可能的 ns1:, ns2: 等（虽然不太可能出现）
    xml_str = re.sub(r"\bns\d+:", "", xml_str)
    xml_str = re.sub(r":ns\d+\b", "", xml_str)
    return xml_str.encode("utf-8")


def _next_rid(existing_ids):
    """生成下一个未被占用的 rId 序号"""
    nums = [
        int(rid[3:])
        for rid in existing_ids
        if rid.startswith("rId") and rid[3:].isdigit()
    ]
    return (max(nums) if nums else 0) + 1


def _update_content_types(root, slide_count, new_tag_parts):
    """更新 [Content_Types].xml 中的 slide Override，并追加 tag 定义"""
    # 先删掉原有的 slide Override 和 notes 相关条目，避免旧顺序影响新 PPT
    for override in list(root.findall(f"{{{CT_NS}}}Override")):
        part = override.get("PartName", "")
        if part.startswith("/ppt/slides/slide"):
            root.remove(override)
        # 删除 notesSlides 和 notesMasters 相关条目
        if "/notesSlides/" in part or "/notesMasters/" in part:
            root.remove(override)

    for idx in range(1, slide_count + 1):
        ET.SubElement(
            root,
            f"{{{CT_NS}}}Override",
            PartName=f"/ppt/slides/slide{idx}.xml",
            ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml",
        )

    for part_name in new_tag_parts:
        ET.SubElement(
            root,
            f"{{{CT_NS}}}Override",
            PartName=f"/{part_name}",
            ContentType="application/vnd.ms-powerpoint.tags+xml",
        )


def _update_presentation_rels(root, slide_count):
    """更新 ppt/_rels/presentation.xml.rels，返回新增 rId 列表"""
    # 统计当前文件使用的 rId，生成下一个可用序号，避免与模板原数据冲突
    existing = [
        rel.get("Id")
        for rel in root.findall(f"{{{PKG_REL_NS}}}Relationship")
        if rel.get("Id")
    ]
    start = _next_rid(existing)

    for rel in list(root.findall(f"{{{PKG_REL_NS}}}Relationship")):
        target = rel.get("Target", "")
        # 删除原有 slide 引用
        if target.startswith("slides/slide"):
            root.remove(rel)
        # 删除 notesSlides 和 notesMasters 引用，避免 PowerPoint 修复提示
        if target.startswith("notesSlides/") or target.startswith("notesMasters/"):
            root.remove(rel)

    new_rel_ids = []
    for idx in range(1, slide_count + 1):
        rid = f"rId{start + idx - 1}"
        new_rel_ids.append(rid)
        ET.SubElement(
            root,
            f"{{{PKG_REL_NS}}}Relationship",
            Id=rid,
            Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide",
            Target=f"slides/slide{idx}.xml",
        )
    return new_rel_ids


def _update_presentation_xml(root, rel_ids):
    """用新的 rId 顺序重建 p:sldIdLst，并删除 notesMasterIdLst"""
    ns = {"p": P_NS}
    sld_id_lst = root.find("p:sldIdLst", ns)
    if sld_id_lst is None:
        sld_id_lst = ET.SubElement(root, f"{{{P_NS}}}sldIdLst")
    else:
        for child in list(sld_id_lst):
            sld_id_lst.remove(child)

    # PowerPoint 要求 slideId 从一个固定值开始，这里沿用 256 起步的策略
    base = 256
    for idx, rid in enumerate(rel_ids):
        attrib = {f"{{{OD_REL_NS}}}id": rid}
        ET.SubElement(
            sld_id_lst,
            f"{{{P_NS}}}sldId",
            attrib,
            id=str(base + idx),
        )

    # 删除 notesMasterIdLst，避免引用已删除的 notesMaster
    notes_master_lst = root.find("p:notesMasterIdLst", ns)
    if notes_master_lst is not None:
        root.remove(notes_master_lst)


def build_from_json(template_path, json_path, output_path):
    """复制原始模板 pptx，并按 JSON 顺序重新组织 slide 文件"""
    data = json.loads(Path(json_path).read_text(encoding="utf-8"))
    pages = data.get("ppt_pages", [])
    if not pages:
        raise ValueError("JSON 中未找到 ppt_pages 内容")

    with tempfile.TemporaryDirectory() as tmpdir:
        temp_copy = Path(tmpdir) / "working.pptx"
        shutil.copyfile(template_path, temp_copy)

        # 将模板 PPT 的所有文件读入内存，便于自由重写
        with zipfile.ZipFile(template_path, "r") as tmpl_zip:
            file_bytes = {name: tmpl_zip.read(name) for name in tmpl_zip.namelist()}

        slide_map = {}
        slide_rel_map = {}
        for name in file_bytes:
            match = SLIDE_RE.fullmatch(name)
            if match:
                slide_map[int(match.group(1))] = file_bytes[name]
            elif name.startswith("ppt/slides/_rels/slide") and name.endswith(
                ".xml.rels"
            ):
                num = int(re.search(r"slide(\d+)\.xml\.rels", name).group(1))
                slide_rel_map[num] = file_bytes[name]

        slide_count = len(slide_map)
        if slide_count == 0:
            raise ValueError("模板中未找到任何 slide 文件")

        tag_nums = [
            int(m.group(1))
            for name in file_bytes
            if (m := TAG_RE.fullmatch(name)) is not None
        ]
        next_tag_num = max(tag_nums) if tag_nums else 0
        extra_tag_parts = []
        extra_tag_files = {}

        # 根据 JSON 记录应使用的模板页编号，同时打印处理进度
        selected_slides = []
        for idx, page in enumerate(pages, start=1):
            tmpl_num = page.get("template_page_num")
            page_type = page.get("page_type", "未知版式")
            if tmpl_num is None:
                raise ValueError(f"第{idx}条缺少 template_page_num")
            if tmpl_num not in slide_map:
                raise ValueError(
                    f"模板中不存在第{tmpl_num}页（来自第{idx}条 {page_type}）"
                )
            selected_slides.append((tmpl_num, page_type))
            print(f"✅ 生成第{idx}页：{page_type}（模板第{tmpl_num}页）")

        slide_total = len(selected_slides)

        pres_rels = ET.fromstring(file_bytes["ppt/_rels/presentation.xml.rels"])
        new_rel_ids = _update_presentation_rels(pres_rels, slide_total)

        pres_xml = ET.fromstring(file_bytes["ppt/presentation.xml"])
        _update_presentation_xml(pres_xml, new_rel_ids)

        def clone_tags(rel_bytes):
            """处理 slide rels：复制 tag 关系并删除 notesSlide 引用"""
            nonlocal next_tag_num
            if not rel_bytes:
                return rel_bytes
            rel_tree = ET.fromstring(rel_bytes)
            for rel in list(rel_tree.findall(f"{{{PKG_REL_NS}}}Relationship")):
                target = rel.get("Target", "")
                rel_type = rel.get("Type", "")

                # 删除 notesSlide 引用，避免 PowerPoint 修复提示
                if "notesSlide" in target or "notesSlide" in rel_type:
                    rel_tree.remove(rel)
                    continue

                # 处理 tag 关系
                if (
                    rel_type
                    == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags"
                ):
                    canonical = posixpath.normpath(posixpath.join("ppt/slides", target))
                    if canonical not in file_bytes:
                        continue
                    # tag 关系在 PPT 中要求唯一，因此为每条关系生成新的 tag 文件
                    next_tag_num += 1
                    new_part = f"ppt/tags/tag{next_tag_num}.xml"
                    rel.set("Target", posixpath.relpath(new_part, "ppt/slides"))
                    extra_tag_parts.append(new_part)
                    extra_tag_files[new_part] = file_bytes[canonical]
            return _clean_rels_namespace(
                ET.tostring(rel_tree, encoding="utf-8", xml_declaration=True)
            )

        prepared_rel_bytes = {}
        for idx, (tmpl_num, _) in enumerate(selected_slides, start=1):
            if tmpl_num in slide_rel_map:
                prepared_rel_bytes[idx] = clone_tags(slide_rel_map[tmpl_num])

        content_types = ET.fromstring(file_bytes["[Content_Types].xml"])
        _update_content_types(content_types, slide_total, extra_tag_parts)

        with zipfile.ZipFile(output_path, "w") as out_zip:
            # 1. 先写入所有与 slide 无关的原始文件（主题、字体、媒体等）
            for name, data in file_bytes.items():
                if name.startswith("ppt/slides/slide"):
                    continue
                if name.startswith("ppt/slides/_rels/slide"):
                    continue
                # 跳过 notesSlides 和 notesMasters 相关文件，避免 PowerPoint 修复提示
                if "notesSlides" in name or "notesMasters" in name:
                    continue
                if name == "[Content_Types].xml":
                    out_zip.writestr(
                        name,
                        ET.tostring(
                            content_types, encoding="utf-8", xml_declaration=True
                        ),
                    )
                elif name == "ppt/_rels/presentation.xml.rels":
                    out_zip.writestr(
                        name,
                        _clean_rels_namespace(
                            ET.tostring(
                                pres_rels, encoding="utf-8", xml_declaration=True
                            )
                        ),
                    )
                elif name == "ppt/presentation.xml":
                    out_zip.writestr(
                        name,
                        ET.tostring(pres_xml, encoding="utf-8", xml_declaration=True),
                    )
                else:
                    out_zip.writestr(name, data)

            # 2. 把 JSON 指定顺序的 slide 与关系文件依次写入
            for idx, (tmpl_num, _) in enumerate(selected_slides, start=1):
                slide_name = f"ppt/slides/slide{idx}.xml"
                rel_name = f"ppt/slides/_rels/slide{idx}.xml.rels"
                out_zip.writestr(slide_name, slide_map[tmpl_num])
                rel_bytes = prepared_rel_bytes.get(idx)
                if rel_bytes:
                    out_zip.writestr(rel_name, rel_bytes)

            # 3. 写入为 tag 关系复制的新文件，确保 PPT 打开不会再修复
            for name, data in extra_tag_files.items():
                out_zip.writestr(name, data)

    print(f"\n🎉 新 PPT 输出完成：{output_path}")


def _create_run_dir(base_dir: Path = Path("temp")) -> Path:
    """创建带时间戳的 run 目录，GUI 可据此收集 PPT 与调试文件。"""
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    suffix = secrets.token_hex(2)
    run_dir = base_dir / f"slide-{timestamp}-{suffix}"
    run_dir.mkdir(parents=True, exist_ok=True)
    return run_dir


def _iter_shapes(shapes):
    """递归遍历幻灯片，包含组合内部的形状。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _iter_shapes_with_path(shapes, parent_path=None):
    """递归遍历幻灯片，返回 (shape, 路径)。"""
    for idx, shape in enumerate(shapes, start=1):
        name = (shape.name or "").strip() or f"元素{idx}"
        path = (*parent_path, name) if parent_path else (name,)
        yield shape, path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_with_path(shape.shapes, path)


def _is_picture_shape(shape):
    """判断形状是否充当图片占位符。"""
    name = shape.name or ""
    # 根据名称判断（兼容手动命名的图片区域）
    if "图片区" in name or name.startswith("图片"):
        return True
    # 优先根据名称判断是否为图片（兼容 FREEFORM 等非标准图片占位符）
    if any(kw in name for kw in ("配图", "图片区", "插图", "照片")):
        return True
    # 标准 PICTURE 类型
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    # 检查是否有图片填充（FREEFORM 类型的图片占位符）
    try:
        if hasattr(shape, "fill") and shape.fill.type == MSO_FILL_TYPE.PICTURE:
            return True
    except Exception as e:
        # 某些形状的 fill 属性访问可能失败，记录调试信息
        print(f"⚠️  检查形状 fill 属性时出错（{shape.name if hasattr(shape, 'name') else '未知'}）：{e}")
    # 占位符类型 18 = 图片占位符
    if shape.is_placeholder:
        try:
            return shape.placeholder_format.type == 18
        except ValueError:
            return False
    return False


def _is_placeholder_shape(shape):
    """识别普通文本占位符，用于目录“文字文字”这类内容。"""
    if not shape.has_text_frame:
        return False
    text = (shape.text_frame.text or "").strip()
    name = shape.name or ""
    if name.startswith("文本框"):
        return True
    return any(keyword in text for keyword in PLACEHOLDER_KEYWORDS)


def _detect_prefix(slide):
    counter = {}
    for _, path in _iter_shapes_with_path(slide.shapes):
        for segment in path:
            prefix = _extract_prefix(segment)
            if prefix:
                counter[prefix] = counter.get(prefix, 0) + 1
    if not counter:
        return None
    return max(counter.items(), key=lambda item: item[1])[0]


def _extract_prefix(name):
    if "_" not in name:
        return None
    prefix = name.split("_", 1)[0]
    return prefix if "页" in prefix else None


def _normalize_path(path, page_prefix):
    start_idx = 0
    if page_prefix:
        start_idx = -1
        for idx, segment in enumerate(path):
            if page_prefix in segment:
                start_idx = idx
                break
        if start_idx == -1:
            return []
    trimmed = []
    for segment in path[start_idx:]:
        seg = _clean_segment(segment, page_prefix)
        if not seg:
            continue
        trimmed.append(seg)
    return trimmed


def _clean_segment(segment, page_prefix):
    seg = segment.strip()
    if not seg:
        return ""
    if page_prefix and seg.startswith(page_prefix + "_"):
        seg = seg[len(page_prefix) + 1 :]
    for keyword in IGNORE_KEYWORDS:
        if keyword in seg:
            return ""
    return seg


def _flatten_content(content):
    """
    展平 content 结构，返回 (path -> value, path -> type, path -> group_path) 三个映射。

    新版格式：{"字段名": {"type": "text/image", "value": "...", "group_path": "..."}}
    旧版格式：{"字段名": "..."} 或嵌套结构
    """
    value_mapping = {}
    type_mapping = {}
    group_mapping = {}  # 记录每个字段的 group_path

    def walk(node, path):
        if isinstance(node, dict):
            # 检查是否是新版格式（有 type 和 value 字段）
            if "type" in node and "value" in node:
                value_mapping[path] = node.get("value", "")
                type_mapping[path] = node.get("type", "text")
                if node.get("group_path"):
                    group_mapping[path] = node.get("group_path")
            else:
                # 嵌套结构，继续遍历
                for key, value in node.items():
                    walk(value, path + (key,))
        else:
            value_mapping[path] = node

    walk(content or {}, tuple())
    return value_mapping, type_mapping, group_mapping


def _shape_aliases(name):
    """为形状名称生成多种别名，提升匹配成功率。"""
    aliases = set()
    clean = name.strip()
    if not clean:
        return aliases

    aliases.update({clean, clean.replace(" ", "")})

    def _add_parts(separator):
        if separator in clean:
            parts = [p for p in clean.split(separator) if p]
            for part in parts:
                aliases.add(part)
                aliases.add(part.replace(" ", ""))

    _add_parts("_")
    _add_parts("-")

    extra = set()
    for alias in aliases:
        for suf in SUFFIXES:
            if alias.endswith(suf):
                trimmed = alias[: -len(suf)]
                extra.add(trimmed)
                extra.add(trimmed.replace(" ", ""))
    aliases.update(extra)
    return aliases


def _shape_tags(shape):
    """提取形状的描述/标题作为标签，便于弱化对 name 的依赖。"""
    tags = set()
    try:
        for node in shape.element.xpath(".//p:cNvPr", namespaces={"p": P_NS}):
            for attr in ("descr", "title"):
                val = (node.get(attr) or "").strip()
                if not val:
                    continue
                tags.add(val)
                tags.update(_shape_aliases(val))
    except Exception:
        return set()
    return tags


def _candidate_keys(key):
    """给定 JSON 中的键名，生成若干匹配候选。"""
    key = key.strip()
    variants = [key, key.replace(" ", "")]
    for suf in SUFFIXES:
        if key.endswith(suf):
            variants.append(key[: -len(suf)])
            variants.append(key[: -len(suf)].replace(" ", ""))
    seen = set()
    result = []
    for variant in variants:
        if variant and variant not in seen:
            result.append(variant)
            seen.add(variant)
    return result


def _take_shape(pool, key, used: Set[int]):
    """从映射中取出首个未使用的形状。"""
    shapes = pool.get(key) or []
    for shape in shapes:
        if id(shape) not in used:
            return shape
    return None


def _copy_run_format(source_run, target_run):
    """复制 run 的格式属性。"""
    try:
        if source_run.font.size:
            target_run.font.size = source_run.font.size
        if source_run.font.bold is not None:
            target_run.font.bold = source_run.font.bold
        if source_run.font.italic is not None:
            target_run.font.italic = source_run.font.italic
        if source_run.font.color.type is not None:
            target_run.font.color.rgb = source_run.font.color.rgb
        if source_run.font.name:
            target_run.font.name = source_run.font.name
    except Exception:
        pass


def _copy_para_format(source_para, target_para):
    """复制段落的格式属性（缩进、对齐、行距等）。"""
    try:
        target_para.alignment = source_para.alignment
        target_para.level = source_para.level
        if source_para.line_spacing:
            target_para.line_spacing = source_para.line_spacing
        if source_para.space_before:
            target_para.space_before = source_para.space_before
        if source_para.space_after:
            target_para.space_after = source_para.space_after
    except Exception:
        pass


def _set_shape_text(shape, text):
    """填充文本时尽量保持原有格式。"""
    if not shape.has_text_frame:
        return

    text = "" if text is None else str(text)
    lines = text.split("\n")
    tf = shape.text_frame

    # 先禁用自动调整大小，防止填充文本时形状被拉伸
    # 必须在修改文本之前设置，否则形状可能已经被拉伸
    tf.auto_size = MSO_AUTO_SIZE.NONE

    if not tf.paragraphs:
        tf.add_paragraph()

    # 保存第一个段落和第一个 run 的格式作为模板
    template_para = tf.paragraphs[0] if tf.paragraphs else None
    template_run = (
        template_para.runs[0] if template_para and template_para.runs else None
    )

    for idx, line in enumerate(lines):
        if idx < len(tf.paragraphs):
            para = tf.paragraphs[idx]
        else:
            # 新建段落并复制格式
            para = tf.add_paragraph()
            if template_para:
                _copy_para_format(template_para, para)

        if para.runs:
            para.runs[0].text = line
            for run in para.runs[1:]:
                run.text = ""
        else:
            # 没有 run 时创建一个并复制格式
            run = para.add_run()
            run.text = line
            if template_run:
                _copy_run_format(template_run, run)

    # 清理多余段落
    for idx in range(len(lines), len(tf.paragraphs)):
        for run in tf.paragraphs[idx].runs:
            run.text = ""


def _safe_remove_shape(shape):
    element_parent = shape.element.getparent()
    if element_parent is not None:
        element_parent.remove(shape.element)


def _replace_picture(slide, shape, image_path):
    """将图片占位符替换为本地图片，保持位置大小比率。

    如果没有提供图片路径或图片文件不存在，删除该形状。
    """
    if not image_path:
        # 没有提供图片，删除形状
        print(f"ℹ️  图片位置 [{shape.name}] 未提供图片，删除形状")
        _safe_remove_shape(shape)
        return

    image_path = Path(image_path)
    if not image_path.is_file():
        # 图片文件不存在，删除形状
        print(f"⚠️  图片文件不可用：{image_path}，删除形状")
        _safe_remove_shape(shape)
        return

    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    name = shape.name

    try:
        with Image.open(image_path) as img:
            img_w, img_h = img.size
    except Exception as e:
        print(f"⚠️  无法读取图片 {image_path}：{e}，删除形状")
        _safe_remove_shape(shape)
        return

    # 防止除零错误
    if img_h == 0 or height == 0:
        print(f"⚠️  图片或形状高度为 0，无法计算比例，删除形状")
        _safe_remove_shape(shape)
        return

    img_ratio = img_w / img_h
    box_ratio = width / height
    if img_ratio > box_ratio:
        new_width = width
        new_height = width / img_ratio if img_ratio != 0 else height
    else:
        new_height = height
        new_width = height * img_ratio

    new_left = int(left + (width - new_width) / 2)
    new_top = int(top + (height - new_height) / 2)
    new_width = int(new_width)
    new_height = int(new_height)

    parent_shapes = getattr(shape, "_parent", None)
    _safe_remove_shape(shape)
    if parent_shapes is None:
        # 某些占位属于组合但 XML 已被剥离，此时无法安全移除，直接保留原父层
        parent_shapes = slide.shapes

    if parent_shapes is None or not hasattr(parent_shapes, "add_picture"):
        parent_shapes = slide.shapes

    new_pic = parent_shapes.add_picture(
        str(image_path), new_left, new_top, width=new_width, height=new_height
    )
    new_pic.name = name


def _fill_slide(slide, page_content, slide_width):
    """根据 JSON 内容把文本和图片写入对应区域。"""
    prefix = _detect_prefix(slide)
    content_map, type_map, group_map = _flatten_content(page_content)
    all_shapes = list(_iter_shapes(slide.shapes))
    shapes_by_name = {}
    shapes_by_exact = {}
    shapes_by_tag = {}
    text_placeholders = []
    used_shapes: Set[int] = set()
    used_text_shapes: Set[int] = set()
    used_picture_shapes: Set[int] = set()
    picture_placeholders = []

    # 收集空内容字段所属的 GROUP（用于删除绑定的装饰元素）
    # 支持两种方式：1. 新格式 JSON 中的 group_path 字段  2. 从 PPT 形状路径推断
    empty_groups = set()
    for path, value in content_map.items():
        if not value and path in group_map:
            # 该字段为空，记录其 GROUP 路径（新格式）
            empty_groups.add(group_map[path])

    # 构建形状名称到 GROUP 路径的映射（用于删除空内容的绑定元素）
    # 注意：使用形状名称作为 key，因为不同遍历中的 shape 对象 id 可能不同
    shape_name_to_group_path = {}
    for shape, raw_path in _iter_shapes_with_path(slide.shapes):
        if len(raw_path) > 1 and shape.name:
            # 记录形状的直接父级 GROUP 路径
            shape_name_to_group_path.setdefault(shape.name, []).append("/".join(raw_path[:-1]))

    for shape in all_shapes:
        if shape.name:
            shapes_by_exact.setdefault(shape.name, []).append(shape)
            for alias in _shape_aliases(shape.name):
                shapes_by_name.setdefault(alias, []).append(shape)
            if _is_placeholder_shape(shape):
                text_placeholders.append(shape)
            if _is_picture_shape(shape):
                picture_placeholders.append(shape)
        for tag in _shape_tags(shape):
            shapes_by_tag.setdefault(tag, []).append(shape)

    def _pop_placeholder(pool, used):
        while pool:
            candidate = pool.pop(0)
            if id(candidate) not in used:
                return candidate
        return None

    for shape, raw_path in _iter_shapes_with_path(slide.shapes):
        label_path = _normalize_path(raw_path, prefix)
        if not label_path:
            continue
        key = tuple(label_path)
        if key not in content_map:
            continue

        # 跳过形状名称本身包含忽略关键词的形状（装饰性元素）
        shape_name = shape.name or ""
        if any(kw in shape_name for kw in IGNORE_KEYWORDS):
            continue

        value = content_map[key]
        field_type = type_map.get(key)  # 获取字段类型

        # 优先根据 JSON 中的 type 字段判断，否则根据形状类型判断
        is_image = field_type == "image" if field_type else _is_picture_shape(shape)

        if is_image:
            _replace_picture(slide, shape, value)
            used_picture_shapes.add(id(shape))
            used_shapes.add(id(shape))
        elif shape.has_text_frame:
            if value:
                _set_shape_text(shape, value)
                used_text_shapes.add(id(shape))
                used_shapes.add(id(shape))
            else:
                # 字段为空：清空文本框，并记录所在 GROUP 以便删除绑定元素
                shape.text_frame.clear()
                used_text_shapes.add(id(shape))
                used_shapes.add(id(shape))
                # 从 raw_path 推断 GROUP 路径（如果形状在 GROUP 内）
                # raw_path 格式如 ('GROUP名', '子GROUP名', '形状名')
                if len(raw_path) > 1:
                    # 取形状的直接父级 GROUP 路径
                    group_path = "/".join(raw_path[:-1])
                    empty_groups.add(group_path)
        else:
            continue

    # 兼容新旧 JSON 格式的键名匹配
    for area_name, raw_value in page_content.items():
        # 新版格式：值是 dict，包含 type/hint/value 等字段
        # 旧版格式：值直接是字符串
        field_type = None  # 从 JSON 获取的字段类型
        if isinstance(raw_value, dict):
            # 检查是否是新版格式（有 type 和 value 字段）
            if "type" in raw_value and "value" in raw_value:
                value = raw_value.get("value", "")
                field_type = raw_value.get("type", "text")  # text 或 image
            else:
                # 嵌套的旧版结构，跳过让 _flatten_content 处理
                continue
        else:
            value = raw_value

        # 首先尝试精确匹配形状名称（向导模式命名后的形状）
        shape = _take_shape(shapes_by_exact, area_name, used_shapes)

        # 如果精确匹配失败，尝试别名匹配
        if not shape:
            for candidate in _candidate_keys(area_name):
                shape = _take_shape(shapes_by_name, candidate, used_shapes)
                if shape:
                    break

        # 标签（alt text/title）匹配，弱化对 name 的依赖
        if not shape:
            tag = _take_shape(shapes_by_tag, area_name, used_shapes)
            if tag:
                shape = tag
        if not shape:
            for candidate in _candidate_keys(area_name):
                shape = _take_shape(shapes_by_tag, candidate, used_shapes)
                if shape:
                    break

        # 尝试手动映射
        if not shape and area_name in MANUAL_NAME_MAP:
            for exact in MANUAL_NAME_MAP[area_name]:
                shape = _take_shape(shapes_by_exact, exact, used_shapes)
                if shape:
                    break
        if not shape and "字幕" in area_name:
            for exact in MANUAL_NAME_MAP.get("字幕", []):
                shape = _take_shape(shapes_by_exact, exact, used_shapes)
                if shape:
                    break

        # 尝试占位符匹配
        if not shape and any(keyword in area_name for keyword in ("内容", "字幕")):
            shape = _pop_placeholder(text_placeholders, used_shapes)
        if not shape and any(keyword in area_name for keyword in ("图片区", "图片")):
            shape = _pop_placeholder(picture_placeholders, used_shapes)

        if not shape:
            print(f"⚠️  找不到名为「{area_name}」的形状，内容已忽略。")
            continue

        # 跳过已处理的形状
        if id(shape) in used_shapes:
            continue

        # 根据 JSON 中的 type 字段或形状类型判断处理方式
        is_image = field_type == "image" if field_type else _is_picture_shape(shape)

        if is_image:
            # 图片类型：没有值时删除形状
            _replace_picture(slide, shape, value)
            used_picture_shapes.add(id(shape))
            used_shapes.add(id(shape))
            # 如果图片为空，记录所在 GROUP
            if not value and shape.name in shape_name_to_group_path:
                for gp in shape_name_to_group_path[shape.name]:
                    empty_groups.add(gp)
        elif shape.has_text_frame:
            # 文本类型：没有值时清空文本框
            if value:
                _set_shape_text(shape, value)
            else:
                shape.text_frame.clear()
                # 记录空文本框所在 GROUP
                if shape.name in shape_name_to_group_path:
                    for gp in shape_name_to_group_path[shape.name]:
                        empty_groups.add(gp)
            used_text_shapes.add(id(shape))
            used_shapes.add(id(shape))
        else:
            print(f"⚠️  形状「{area_name}」既不是文本也不是图片，跳过。")

    _clear_default_subtitles(all_shapes)
    _apply_layout_rules(all_shapes, slide_width)

    # 处理空内容 GROUP：删除属于空 GROUP 的所有形状
    if empty_groups:
        _delete_empty_group_shapes(slide, empty_groups)


def _delete_empty_group_shapes(slide, empty_groups):
    """
    删除属于空内容 GROUP 的所有形状。

    当某个内容字段为空时，应该删除与其绑定的整个 GROUP（包括装饰元素）。

    Args:
        slide: 幻灯片对象
        empty_groups: 需要删除的 GROUP 路径集合
    """

    def find_and_delete_groups(shapes, parent_path=""):
        """递归查找并删除匹配的 GROUP"""
        groups_to_delete = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_name = shape.name or ""
                current_path = (
                    f"{parent_path}/{group_name}" if parent_path else group_name
                )

                # 检查当前 GROUP 是否应该删除：
                # 1. 完全匹配 empty_groups 中的路径
                # 2. 是 empty_groups 中某个路径的子 GROUP（如 "组合1/子组合" 是 "组合1" 的子级）
                # 不删除外层 GROUP
                should_delete = any(
                    current_path == eg or current_path.startswith(eg + "/")
                    for eg in empty_groups
                )

                if should_delete:
                    groups_to_delete.append(shape)
                    # 删除 GROUP 会自动删除其所有子元素，无需递归
                else:
                    # 递归检查子 GROUP
                    find_and_delete_groups(shape.shapes, current_path)

        # 删除标记的 GROUP（删除 GROUP 会自动删除其所有子元素）
        for group in groups_to_delete:
            try:
                sp = group._element
                sp.getparent().remove(sp)
                print(f"🗑️  已删除空内容 GROUP：{group.name}")
            except Exception as e:
                print(f"⚠️  删除 GROUP 失败：{group.name}, 错误：{e}")

    find_and_delete_groups(slide.shapes)


def _clear_default_subtitles(shapes):
    """清除未被覆盖的字幕或默认说明文字。"""
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if text in SUBTITLE_TEXTS or any(
            keyword in text for keyword in PLACEHOLDER_KEYWORDS
        ):
            shape.text_frame.clear()


def _apply_layout_rules(shapes, slide_width):
    """根据文本长度自动调整标题条与背景。"""
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        name = shape.name or ""
        if not any(keyword in name for keyword in EXPANDABLE_KEYWORDS):
            continue
        _adjust_text_shape(shape, shapes, slide_width)


def _adjust_text_shape(text_shape, shapes, slide_width):
    text_width = _estimate_text_width(text_shape)
    if text_width <= 0:
        return

    limit = _find_right_limit(text_shape, shapes, slide_width) - H_PADDING
    available = max(text_shape.width, limit - text_shape.left)
    target_width = min(max(text_shape.width, text_width + H_PADDING), available)

    if target_width > text_shape.width:
        target_width = int(target_width)
        # 保存原始位置和高度（占位符的这些值可能来自 layout）
        original_left = text_shape.left
        original_top = text_shape.top
        original_height = text_shape.height
        # 修改宽度（这可能会创建新的 xfrm 元素，但只有 ext 没有 off）
        text_shape.width = target_width
        # 强制设置位置和高度，确保 xfrm 元素完整（包含 off 和 ext）
        # 即使值相同也要设置，因为 python-pptx 在设置 width 时可能没有创建 off 元素
        text_shape.left = original_left
        text_shape.top = original_top
        if original_height > 0:
            text_shape.height = original_height
        bg = _find_background_shape(text_shape, shapes)
        if bg:
            bg.width = int(max(bg.width, target_width + H_PADDING))
            bg.left = min(bg.left, text_shape.left)
    else:
        shrink_ratio = available / text_width if text_width else 1
        if shrink_ratio < 1:
            _shrink_font(text_shape, shrink_ratio)


def _find_background_shape(text_shape, shapes):
    top = text_shape.top
    bottom = text_shape.top + text_shape.height
    candidate = None
    for shape in shapes:
        if shape is text_shape or shape.has_text_frame or _is_picture_shape(shape):
            continue
        overlap = min(bottom, shape.top + shape.height) - max(top, shape.top)
        if overlap <= 0:
            continue
        ratio = overlap / max(1, text_shape.height)
        if ratio < 0.6:
            continue
        if candidate is None or shape.width > candidate.width:
            candidate = shape
    return candidate


def _find_right_limit(text_shape, shapes, slide_width):
    limit = slide_width
    top = text_shape.top
    bottom = text_shape.top + text_shape.height
    for shape in shapes:
        if shape is text_shape:
            continue
        other_top = shape.top
        other_bottom = shape.top + shape.height
        overlap = min(bottom, other_bottom) - max(top, other_top)
        if overlap <= 0:
            continue
        if shape.left > text_shape.left:
            limit = min(limit, shape.left)
    return limit


def _estimate_text_width(shape):
    if not shape.has_text_frame:
        return 0

    max_line = 0
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs) or para.text or ""
        if not line:
            continue
        font_size = None
        for run in para.runs:
            if run.font.size:
                font_size = run.font.size.pt
                break
        if font_size is None:
            font_size = 28
        width_factor = sum(0.55 if ord(ch) < 128 else 1 for ch in line)
        line_width = width_factor * font_size * EMU_PER_PT
        max_line = max(max_line, line_width)
    return max(max_line, shape.width)


def _shrink_font(shape, ratio):
    ratio = max(ratio, 0.6)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                run.font.size = Pt(run.font.size.pt * ratio)


def _extract_connectors(pptx_path: Path):
    connectors = {}
    with zipfile.ZipFile(pptx_path, "r") as zf:
        for name in zf.namelist():
            match = SLIDE_RE.fullmatch(name)
            if not match:
                continue
            slide_idx = int(match.group(1))
            root = etree.fromstring(zf.read(name))
            nodes = root.xpath(".//p:cxnSp", namespaces=NSMAP)
            if nodes:
                connectors[slide_idx] = [etree.tostring(node) for node in nodes]
    return connectors


def _restore_connectors(pptx_path: Path, connectors):
    if not connectors:
        return
    with zipfile.ZipFile(pptx_path, "r") as src:
        entries = {name: src.read(name) for name in src.namelist()}
    modified = False
    for name, data in list(entries.items()):
        match = SLIDE_RE.fullmatch(name)
        if not match:
            continue
        slide_idx = int(match.group(1))
        snippets = connectors.get(slide_idx)
        if not snippets:
            continue
        root = etree.fromstring(data)
        sp_tree = root.find(".//p:spTree", namespaces=NSMAP)
        if sp_tree is None:
            continue
        for node in sp_tree.findall("p:cxnSp", namespaces=NSMAP):
            sp_tree.remove(node)
        for snippet in snippets:
            sp_tree.append(etree.fromstring(snippet))
        entries[name] = etree.tostring(root, encoding="utf-8", xml_declaration=True)
        modified = True
    if not modified:
        return
    with zipfile.ZipFile(pptx_path, "w") as dst:
        for name, data in entries.items():
            dst.writestr(name, data)


def render_slides(
    template_path: Path,
    config: Dict,
    output_name: str,
    run_dir: Optional[Path] = None,
) -> Dict:
    """渲染入口，供 GUI/CLI 复用，返回 PPT 路径和 run 目录信息。"""
    pages = config.get("ppt_pages", [])
    if not pages:
        raise ValueError("JSON 数据中没有 ppt_pages 内容。")

    run_dir = run_dir or _create_run_dir()
    run_dir.mkdir(parents=True, exist_ok=True)
    output_path = run_dir / output_name

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        temp_ppt = Path(tmp.name)

    try:
        tmp_json = run_dir / "config.json"
        tmp_json.write_text(
            json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8"
        )

        build_from_json(template_path, tmp_json, temp_ppt)
        connector_snapshots = _extract_connectors(temp_ppt)
        prs = Presentation(temp_ppt)
        if len(prs.slides) != len(pages):
            raise RuntimeError("生成的幻灯片数量与 JSON 不匹配，无法填充。")

        slide_width = prs.slide_width
        for slide, page in zip(prs.slides, pages):
            _fill_slide(slide, page.get("content", {}), slide_width)

        prs.save(output_path)
        _restore_connectors(output_path, connector_snapshots)
    finally:
        temp_ppt.unlink(missing_ok=True)

    return {"output_path": output_path, "run_dir": run_dir, "slides": len(pages)}


def main():
    parser = argparse.ArgumentParser(description="读取 JSON 并填充模板生成 PPT")
    parser.add_argument("--template", required=True, help="模板 PPTX 路径")
    parser.add_argument("--json", required=True, help="描述内容的 JSON 文件")
    parser.add_argument(
        "--output", default="final_output.pptx", help="输出 PPTX 文件名或路径"
    )
    parser.add_argument(
        "--run-dir", default=None, help="输出 run 目录（默认 temp/run-...）"
    )
    args = parser.parse_args()

    config = json.loads(Path(args.json).read_text(encoding="utf-8"))
    run_dir = Path(args.run_dir) if args.run_dir else None
    result = render_slides(Path(args.template), config, Path(args.output).name, run_dir)
    final_path = result["output_path"]
    if Path(args.output).is_absolute():
        Path(args.output).parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(final_path, Path(args.output))
        print(f"📄 另存为：{args.output}")
    print(f"🎯 已根据内容生成 PPT：{final_path}")


if __name__ == "__main__":
    main()
