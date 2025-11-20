"""
从模板 PPT 中提取可供大模型参考的 JSON 结构。

用法示例：
python scripts/export_template_structure.py \
    --template template/template.pptx \
    --output template/exported_template.json \
    --mode semantic \
    --include 1,2,3,4,8,12,15,16,17,18,21,26,27,28 \
    --ai-enrich --llm-provider deepseek --llm-model deepseek-chat
"""

from __future__ import annotations

import argparse
import json
import math
import os
import re
from collections import Counter, OrderedDict
from pathlib import Path
from typing import Dict, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

GENERIC_NAME_PATTERN = re.compile(
    r"^(图片|文本框|矩形|圆角|任意|椭圆|线条|组合|对象|table|textbox|picture|group|背景|subtitle|caption)",
    re.IGNORECASE,
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="从 PPT 模板导出符合 template.json 结构的描述文件。"
    )
    parser.add_argument("--template", required=True, help="PPTX 模板路径")
    parser.add_argument("--output", required=True, help="导出 JSON 的输出路径")
    parser.add_argument(
        "--mode",
        choices=("semantic", "text"),
        default="semantic",
        help=(
            "semantic：仅导出命名规范（含“xx区”等）的元素；"
            "text：导出所有可编辑文本框（忽略图片/背景）。"
        ),
    )
    parser.add_argument(
        "--include",
        help="可选，逗号分隔的页码列表，仅导出这些幻灯片，例如：1,2,4",
    )
    parser.add_argument(
        "--ai-enrich",
        action="store_true",
        help="使用 AI 自动填充 hint、required、max_chars 和 notes 字段",
    )
    parser.add_argument(
        "--llm-provider",
        choices=("deepseek", "local", "qwen"),
        default="deepseek",
        help="LLM 提供商（仅在 --ai-enrich 时有效）",
    )
    parser.add_argument(
        "--llm-model",
        help="LLM 模型名称（仅在 --ai-enrich 时有效）",
    )
    parser.add_argument(
        "--llm-base-url",
        help="LLM 服务器地址（仅在 --ai-enrich 时有效）",
    )
    return parser.parse_args()


def is_meaningful_name(name: str) -> bool:
    name = name.strip()
    if not name:
        return False
    if GENERIC_NAME_PATTERN.match(name):
        return False
    if "_" in name:
        return True
    if re.search(r"\d", name):
        return False
    return True


def sanitize_name(name: str, fallback: str) -> str:
    cleaned = name.strip()
    return cleaned or fallback


def flatten_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    parts = [para.text for para in shape.text_frame.paragraphs]
    return "\n".join(part for part in parts if part)


def estimate_max_chars(sample: str) -> int:
    sample = re.sub(r"\s+", "", sample or "")
    length = len(sample)
    if length == 0:
        return 20
    return max(6, min(60, math.ceil(length * 1.5)))


def should_include_text_shape(
    name: str, text: str, mode: str, context_has_semantics: bool
) -> bool:
    if not text.strip():
        return False
    if mode == "semantic":
        if is_meaningful_name(name):
            return True
        return context_has_semantics
    if text.strip().startswith("字幕"):
        return False
    return True


def should_include_group(name: str, mode: str) -> bool:
    if mode == "text":
        return False
    return is_meaningful_name(name)


def is_image_shape(shape) -> bool:
    return shape.shape_type in {
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.LINKED_PICTURE,
    }


def infer_page_type(slide, fallback: str) -> str:
    candidates: Counter[str] = Counter()

    def register(name: str) -> None:
        name = name.strip()
        if name:
            candidates[name] += 1

    for shape in slide.shapes:
        name = shape.name.strip()
        if not name:
            continue
        meaningful = is_meaningful_name(name)
        if "_" in name and meaningful:
            prefix = name.split("_", 1)[0]
            if "页" in prefix:
                register(prefix)
        if meaningful and "页" in name:
            register(name)
        if meaningful and "多字版" in name:
            register("章节页多字版")
            register("章节页多字版")
        if meaningful and "章节" in name:
            register("章节页")
        if meaningful and name.startswith("图文页"):
            register(name)
        if meaningful and name.startswith("文字页"):
            register(name)

    if candidates:
        best = candidates.most_common()
        top = best[0][1]
        tied = [name for name, cnt in best if cnt == top]
        return max(tied, key=len)

    texts = " ".join(
        flatten_text(shape) for shape in slide.shapes if shape.has_text_frame
    )
    keyword_map = [
        ("目录", "目录页"),
        ("主讲", "主讲人页"),
        ("章节", "章节页"),
        ("过渡", "过渡页"),
        ("图文", "图文页"),
    ]
    for kw, page in keyword_map:
        if kw in texts:
            return page
    return fallback


def add_field(container: OrderedDict, path: Sequence[str], payload: Dict) -> None:
    cursor = container
    for key in path[:-1]:
        cursor = cursor.setdefault(key, OrderedDict())
    cursor[path[-1]] = payload


def collect_fields(slide, mode: str) -> Tuple[OrderedDict, int, int]:
    content: OrderedDict = OrderedDict()
    text_slots = 0
    image_slots = 0
    auto_counters: Dict[str, int] = {}

    def visit(shape, context: Tuple[str, ...], context_has_semantics: bool):
        nonlocal text_slots, image_slots
        name = shape.name.strip()
        has_semantics = context_has_semantics or is_meaningful_name(name)

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            new_context = context
            new_semantics = context_has_semantics
            if should_include_group(name, mode):
                new_context = context + (
                    sanitize_name(name, f"区域{len(context) + 1}"),
                )
                new_semantics = True
            for child in shape.shapes:
                visit(child, new_context, new_semantics)
            return

        if mode == "semantic" and is_image_shape(shape) and is_meaningful_name(name):
            field_name = sanitize_name(name.split("_")[-1], name)
            add_field(
                content,
                context + (field_name,),
                {
                    "type": "image",
                    "hint": f"为“{field_name}”提供图片路径",
                    "required": True,
                    "value": "",
                    "preferred_format": "png/jpg",
                },
            )
            image_slots += 1
            return

        if shape.has_text_frame:
            text = flatten_text(shape).strip()
            if not should_include_text_shape(name, text, mode, context_has_semantics):
                return

            key_base = name if is_meaningful_name(name) else None
            if not key_base:
                prefix = "文本"
                auto_counters[prefix] = auto_counters.get(prefix, 0) + 1
                key_base = f"{prefix}{auto_counters[prefix]}"

            field_path = context + (sanitize_name(key_base, key_base),)
            add_field(
                content,
                field_path,
                {
                    "type": "text",
                    "hint": f"填写“{field_path[-1]}”的内容",
                    "required": True,
                    "value": "",
                    "max_chars": estimate_max_chars(text),
                },
            )
            text_slots += 1

    for shp in slide.shapes:
        visit(shp, tuple(), False)

    return content, text_slots, image_slots


def build_manifest_entry(
    page_num: int, page_type: str, text_slots: int, image_slots: int
) -> Dict:
    return {
        "template_page_num": page_num,
        "page_type": page_type,
        "text_slots": text_slots,
        "image_slots": image_slots,
    }


def export_template_structure(
    template_path: Path, mode: str, include_pages: Optional[Sequence[int]]
) -> Dict:
    prs = Presentation(str(template_path))
    data_manifest: List[Dict] = []
    ppt_pages: List[Dict] = []
    include_set = set(include_pages) if include_pages else None

    for idx, slide in enumerate(prs.slides, start=1):
        if include_set and idx not in include_set:
            continue
        page_type = infer_page_type(slide, fallback=f"模板第{idx}页")
        content, text_slots, image_slots = collect_fields(slide, mode)
        if not content and mode == "semantic":
            continue
        manifest_entry = build_manifest_entry(idx, page_type, text_slots, image_slots)
        page_payload = {
            "page_type": page_type,
            "template_page_num": idx,
            "content": content,
            "meta": {
                "layout": "auto-generated",
                "scene": [],
                "style": "auto",
                "text_slots": text_slots,
                "image_slots": image_slots,
                "notes": (
                    "由 export_template_structure.py 自动提取。"
                    "请根据实际模板补充更丰富的描述。"
                ),
            },
        }
        data_manifest.append(manifest_entry)
        ppt_pages.append(page_payload)

    return {"manifest": data_manifest, "ppt_pages": ppt_pages}


def ai_enrich_template(
    template_data: Dict,
    llm_provider: str = "deepseek",
    llm_model: Optional[str] = None,
    llm_base_url: Optional[str] = None,
) -> Dict:
    """使用 AI 自动填充模板配置中的 hint、required、max_chars 和 notes 字段。

    Args:
        template_data: 从 export_template_structure 生成的模板数据
        llm_provider: LLM 提供商 (deepseek/local/qwen)
        llm_model: LLM 模型名称
        llm_base_url: LLM 服务器地址

    Returns:
        填充后的模板数据
    """
    # Import LLM client
    try:
        from scripts.llm_client import BaseLLM, DeepSeekLLM, LocalLLM, QwenVLLM
    except ImportError:
        # Fallback for when running from web context
        from llm_client import BaseLLM, DeepSeekLLM, LocalLLM, QwenVLLM

    # Initialize LLM
    llm: BaseLLM
    provider = llm_provider.lower()
    if provider == "deepseek":
        llm = DeepSeekLLM(model=llm_model or "deepseek-chat")
    elif provider == "local":
        llm = LocalLLM(model=llm_model, base_url=llm_base_url)
    elif provider == "qwen":
        if not llm_base_url:
            llm_base_url = os.getenv("QWEN_VLLM_BASE_URL")
        if not llm_base_url:
            raise ValueError(
                "Qwen provider 需要提供 --llm-base-url 或设置 QWEN_VLLM_BASE_URL"
            )
        llm = QwenVLLM(base_url=llm_base_url)
    else:
        raise ValueError(f"不支持的 LLM 提供商：{llm_provider}")

    print(f"🤖 使用 {llm_provider} 进行 AI 填充...")

    # Process each page
    enriched_data = template_data.copy()
    enriched_data["ppt_pages"] = []

    for page_idx, page in enumerate(template_data["ppt_pages"], start=1):
        print(
            f"  处理第 {page_idx}/{len(template_data['ppt_pages'])} 页: {page['page_type']}"
        )

        # Build prompt for this page
        page_type = page["page_type"]
        template_page_num = page["template_page_num"]
        fields = page.get("content", {})

        # Create field summary
        field_names = list(fields.keys())
        field_summary = "\n".join([f"  - {name}" for name in field_names])

        prompt = f"""你是一个 PPT 模板配置专家。现在需要为一个 PPT 模板页面填写配置信息。

页面信息：
- 页面类型：{page_type}
- 模板页码：{template_page_num}
- 字段列表：
{field_summary}

请为这个页面生成配置信息，包括：
1. 页面说明 (notes)：简要描述这一页的用途和内容要求（1-2句话）
2. 每个字段的配置：
   - hint：提示大模型该字段应该填写什么内容（简洁明了）
   - required：该字段是否必填（true/false）
   - max_chars：该字段的最大字符数（合理估计）

请以 JSON 格式返回，格式如下：
{{
  "notes": "页面说明文字",
  "fields": {{
    "字段名1": {{
      "hint": "提示文字",
      "required": true,
      "max_chars": 20
    }},
    "字段名2": {{
      "hint": "提示文字",
      "required": false,
      "max_chars": 50
    }}
  }}
}}

注意：
- notes 要简洁明了，帮助大模型理解页面用途
- hint 要具体，说明该字段应该填什么内容
- required 根据字段的重要性判断
- max_chars 要合理，考虑页面布局和内容需求
- 只返回 JSON，不要有其他文字"""

        try:
            # Call LLM
            response = llm.generate(
                messages=[{"role": "user", "content": prompt}],
                temperature=0.3,
            )

            # Parse response
            # Try to extract JSON from response
            json_match = re.search(r"\{[\s\S]*\}", response)
            if json_match:
                ai_config = json.loads(json_match.group())
            else:
                ai_config = json.loads(response)

            # Update page notes
            if "notes" in ai_config:
                page["meta"]["notes"] = ai_config["notes"]

            # Update field configurations
            if "fields" in ai_config:
                for field_name, field_config in ai_config["fields"].items():
                    if field_name in fields:
                        if "hint" in field_config:
                            fields[field_name]["hint"] = field_config["hint"]
                        if "required" in field_config:
                            fields[field_name]["required"] = field_config["required"]
                        if "max_chars" in field_config:
                            fields[field_name]["max_chars"] = field_config["max_chars"]

            print(f"    ✅ 成功填充")

        except Exception as e:
            print(f"    ⚠️  AI 填充失败: {e}")
            print(f"    使用默认配置")

        enriched_data["ppt_pages"].append(page)

    print("✅ AI 填充完成")
    return enriched_data


def main() -> None:
    args = parse_args()
    template_path = Path(args.template)
    if not template_path.exists():
        raise FileNotFoundError(f"模板文件不存在：{template_path}")

    include_pages = None
    if args.include:
        include_pages = [
            int(num.strip()) for num in args.include.split(",") if num.strip()
        ]

    data = export_template_structure(template_path, args.mode, include_pages)

    # AI enrichment if requested
    if args.ai_enrich:
        data = ai_enrich_template(
            template_data=data,
            llm_provider=args.llm_provider,
            llm_model=args.llm_model,
            llm_base_url=args.llm_base_url,
        )

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with output_path.open("w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"已导出 {len(data['ppt_pages'])} 个模板页面 -> {output_path}")


if __name__ == "__main__":
    main()
