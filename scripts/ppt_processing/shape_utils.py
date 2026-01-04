"""形状检测和遍历工具。"""

from typing import Generator, List, Optional, Set, Tuple

from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE

from scripts.ppt_processing.constants import (
    SUFFIXES,
    PLACEHOLDER_KEYWORDS,
    IGNORE_KEYWORDS,
    P_NS,
)


def iter_shapes(shapes) -> Generator:
    """递归遍历幻灯片，包含组合内部的形状。

    Args:
        shapes: 形状集合

    Yields:
        每个形状对象
    """
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def iter_shapes_with_path(
    shapes,
    parent_path: Optional[Tuple[str, ...]] = None
) -> Generator[Tuple, None, None]:
    """递归遍历幻灯片，返回 (shape, 路径)。

    Args:
        shapes: 形状集合
        parent_path: 父级路径元组

    Yields:
        (shape, path) 元组
    """
    for idx, shape in enumerate(shapes, start=1):
        name = (shape.name or "").strip() or f"元素{idx}"
        path = (*parent_path, name) if parent_path else (name,)
        yield shape, path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes_with_path(shape.shapes, path)


def is_picture_shape(shape) -> bool:
    """判断形状是否充当图片占位符。

    Args:
        shape: 形状对象

    Returns:
        是否为图片形状
    """
    name = shape.name or ""

    # 根据名称判断（兼容手动命名的图片区域）
    if "图片区" in name or name.startswith("图片"):
        return True

    # 优先根据名称判断是否为图片
    if any(kw in name for kw in ("配图", "图片区", "插图", "照片")):
        return True

    # 标准 PICTURE 类型
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True

    # 检查是否有图片填充（FREEFORM 类型的图片占位符）
    try:
        if hasattr(shape, "fill") and shape.fill.type == MSO_FILL_TYPE.PICTURE:
            return True
    except Exception as e:
        shape_name = shape.name if hasattr(shape, "name") else "未知"
        print(f"⚠️  检查形状 fill 属性时出错（{shape_name}）：{e}")

    # 占位符类型 18 = 图片占位符
    if shape.is_placeholder:
        try:
            return shape.placeholder_format.type == 18
        except ValueError:
            return False

    return False


def is_placeholder_shape(shape) -> bool:
    """识别普通文本占位符，用于目录"文字文字"这类内容。

    Args:
        shape: 形状对象

    Returns:
        是否为文本占位符
    """
    if not shape.has_text_frame:
        return False
    text = (shape.text_frame.text or "").strip()
    name = shape.name or ""
    if name.startswith("文本框"):
        return True
    return any(keyword in text for keyword in PLACEHOLDER_KEYWORDS)


def shape_aliases(name: str) -> Set[str]:
    """为形状名称生成多种别名，提升匹配成功率。

    Args:
        name: 形状名称

    Returns:
        别名集合
    """
    aliases = set()
    clean = name.strip()
    if not clean:
        return aliases

    aliases.update({clean, clean.replace(" ", "")})

    def _add_parts(separator: str) -> None:
        if separator in clean:
            parts = [p for p in clean.split(separator) if p]
            for part in parts:
                aliases.add(part)
                aliases.add(part.replace(" ", ""))

    _add_parts("_")
    _add_parts("-")

    extra = set()
    for alias in aliases:
        for suf in SUFFIXES:
            if alias.endswith(suf):
                trimmed = alias[: -len(suf)]
                extra.add(trimmed)
                extra.add(trimmed.replace(" ", ""))
    aliases.update(extra)
    return aliases


def shape_tags(shape) -> Set[str]:
    """提取形状的描述/标题作为标签，便于弱化对 name 的依赖。

    Args:
        shape: 形状对象

    Returns:
        标签集合
    """
    tags = set()
    try:
        for node in shape.element.xpath(".//p:cNvPr", namespaces={"p": P_NS}):
            for attr in ("descr", "title"):
                val = (node.get(attr) or "").strip()
                if not val:
                    continue
                tags.add(val)
                tags.update(shape_aliases(val))
    except Exception:
        return set()
    return tags


def detect_prefix(slide) -> Optional[str]:
    """检测幻灯片中形状名称的页面前缀。

    Args:
        slide: 幻灯片对象

    Returns:
        检测到的前缀，如 "第1页"
    """
    counter = {}
    for _, path in iter_shapes_with_path(slide.shapes):
        for segment in path:
            prefix = _extract_prefix(segment)
            if prefix:
                counter[prefix] = counter.get(prefix, 0) + 1
    if not counter:
        return None
    return max(counter.items(), key=lambda item: item[1])[0]


def _extract_prefix(name: str) -> Optional[str]:
    """从形状名称中提取页面前缀。"""
    if "_" not in name:
        return None
    prefix = name.split("_", 1)[0]
    return prefix if "页" in prefix else None


def normalize_path(
    path: Tuple[str, ...],
    page_prefix: Optional[str]
) -> List[str]:
    """规范化形状路径，移除页面前缀和忽略关键词。

    Args:
        path: 原始路径元组
        page_prefix: 页面前缀

    Returns:
        规范化后的路径列表
    """
    start_idx = 0
    if page_prefix:
        start_idx = -1
        for idx, segment in enumerate(path):
            if page_prefix in segment:
                start_idx = idx
                break
        if start_idx == -1:
            return []

    trimmed = []
    for segment in path[start_idx:]:
        seg = _clean_segment(segment, page_prefix)
        if not seg:
            continue
        trimmed.append(seg)
    return trimmed


def _clean_segment(segment: str, page_prefix: Optional[str]) -> str:
    """清理路径段，移除前缀和忽略关键词。"""
    seg = segment.strip()
    if not seg:
        return ""
    if page_prefix and seg.startswith(page_prefix + "_"):
        seg = seg[len(page_prefix) + 1:]
    for keyword in IGNORE_KEYWORDS:
        if keyword in seg:
            return ""
    return seg


def candidate_keys(key: str) -> List[str]:
    """给定 JSON 中的键名，生成若干匹配候选。

    Args:
        key: JSON 键名

    Returns:
        候选键名列表
    """
    key = key.strip()
    variants = [key, key.replace(" ", "")]
    for suf in SUFFIXES:
        if key.endswith(suf):
            variants.append(key[: -len(suf)])
            variants.append(key[: -len(suf)].replace(" ", ""))
    seen = set()
    result = []
    for variant in variants:
        if variant and variant not in seen:
            result.append(variant)
            seen.add(variant)
    return result


def take_shape(pool: dict, key: str, used: Set[int]):
    """从映射中取出首个未使用的形状。

    Args:
        pool: 形状名称到形状列表的映射
        key: 要查找的键名
        used: 已使用形状的 id 集合

    Returns:
        找到的形状，或 None
    """
    shapes = pool.get(key) or []
    for shape in shapes:
        if id(shape) not in used:
            return shape
    return None
