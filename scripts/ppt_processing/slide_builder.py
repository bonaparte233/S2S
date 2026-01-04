"""幻灯片构建和填充核心逻辑。"""

import json
import posixpath
import re
import secrets
import shutil
import tempfile
import zipfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Set, Tuple
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from scripts.ppt_processing.constants import (
    SLIDE_RE,
    TAG_RE,
    PKG_REL_NS,
    IGNORE_KEYWORDS,
    MANUAL_NAME_MAP,
)
from scripts.ppt_processing.xml_utils import (
    clean_rels_namespace,
    update_content_types,
    update_presentation_rels,
    update_presentation_xml,
)
from scripts.ppt_processing.shape_utils import (
    iter_shapes,
    iter_shapes_with_path,
    is_picture_shape,
    is_placeholder_shape,
    shape_aliases,
    shape_tags,
    detect_prefix,
    normalize_path,
    candidate_keys,
    take_shape,
)
from scripts.ppt_processing.text_utils import set_shape_text
from scripts.ppt_processing.image_utils import replace_picture
from scripts.ppt_processing.layout_utils import (
    apply_layout_rules,
    clear_default_subtitles,
)
from scripts.ppt_processing.connector_utils import (
    extract_connectors,
    restore_connectors,
)


def create_run_dir(base_dir: Path = Path("temp")) -> Path:
    """创建带时间戳的 run 目录。

    Args:
        base_dir: 基础目录

    Returns:
        创建的 run 目录路径
    """
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    suffix = secrets.token_hex(2)
    run_dir = base_dir / f"slide-{timestamp}-{suffix}"
    run_dir.mkdir(parents=True, exist_ok=True)
    return run_dir


def flatten_content(content: Optional[Dict]) -> Tuple[Dict, Dict, Dict]:
    """展平 content 结构，返回三个映射。

    新版格式：{"字段名": {"type": "text/image", "value": "...", "group_path": "..."}}
    旧版格式：{"字段名": "..."} 或嵌套结构

    Args:
        content: 内容字典

    Returns:
        (value_mapping, type_mapping, group_mapping) 三元组
    """
    value_mapping = {}
    type_mapping = {}
    group_mapping = {}

    def walk(node, path: Tuple):
        if isinstance(node, dict):
            if "type" in node and "value" in node:
                value_mapping[path] = node.get("value", "")
                type_mapping[path] = node.get("type", "text")
                if node.get("group_path"):
                    group_mapping[path] = node.get("group_path")
            else:
                for key, value in node.items():
                    walk(value, path + (key,))
        else:
            value_mapping[path] = node

    walk(content or {}, tuple())
    return value_mapping, type_mapping, group_mapping


def delete_empty_group_shapes(slide, empty_groups: Set[str]) -> None:
    """删除属于空内容 GROUP 的所有形状。

    Args:
        slide: 幻灯片对象
        empty_groups: 需要删除的 GROUP 路径集合
    """
    def find_and_delete_groups(shapes, parent_path: str = ""):
        groups_to_delete = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_name = shape.name or ""
                current_path = (
                    f"{parent_path}/{group_name}" if parent_path else group_name
                )
                should_delete = any(
                    current_path == eg or current_path.startswith(eg + "/")
                    for eg in empty_groups
                )
                if should_delete:
                    groups_to_delete.append(shape)
                else:
                    find_and_delete_groups(shape.shapes, current_path)

        for group in groups_to_delete:
            try:
                sp = group._element
                sp.getparent().remove(sp)
                print(f"🗑️  已删除空内容 GROUP：{group.name}")
            except Exception as e:
                print(f"⚠️  删除 GROUP 失败：{group.name}, 错误：{e}")

    find_and_delete_groups(slide.shapes)


def fill_slide(slide, page_content: Dict, slide_width: int) -> None:
    """根据 JSON 内容把文本和图片写入对应区域。

    Args:
        slide: 幻灯片对象
        page_content: 页面内容字典
        slide_width: 幻灯片宽度（EMU）
    """
    prefix = detect_prefix(slide)
    content_map, type_map, group_map = flatten_content(page_content)
    all_shapes = list(iter_shapes(slide.shapes))
    shapes_by_name = {}
    shapes_by_exact = {}
    shapes_by_tag = {}
    text_placeholders = []
    used_shapes: Set[int] = set()
    picture_placeholders = []

    # 收集空内容字段所属的 GROUP
    empty_groups = set()
    for path, value in content_map.items():
        if not value and path in group_map:
            empty_groups.add(group_map[path])

    # 构建形状名称到 GROUP 路径的映射
    shape_name_to_group_path = {}
    for shape, raw_path in iter_shapes_with_path(slide.shapes):
        if len(raw_path) > 1 and shape.name:
            shape_name_to_group_path.setdefault(shape.name, []).append(
                "/".join(raw_path[:-1])
            )

    # 构建形状索引
    for shape in all_shapes:
        if shape.name:
            shapes_by_exact.setdefault(shape.name, []).append(shape)
            for alias in shape_aliases(shape.name):
                shapes_by_name.setdefault(alias, []).append(shape)
            if is_placeholder_shape(shape):
                text_placeholders.append(shape)
            if is_picture_shape(shape):
                picture_placeholders.append(shape)
        for tag in shape_tags(shape):
            shapes_by_tag.setdefault(tag, []).append(shape)

    def _pop_placeholder(pool: List, used: Set[int]):
        while pool:
            candidate = pool.pop(0)
            if id(candidate) not in used:
                return candidate
        return None

    # 第一遍：通过路径匹配
    _fill_by_path(
        slide, prefix, content_map, type_map, used_shapes, empty_groups
    )

    # 第二遍：通过键名匹配
    _fill_by_key(
        slide, page_content, shapes_by_exact, shapes_by_name, shapes_by_tag,
        text_placeholders, picture_placeholders, used_shapes,
        shape_name_to_group_path, empty_groups, _pop_placeholder
    )

    clear_default_subtitles(all_shapes)
    apply_layout_rules(all_shapes, slide_width)

    if empty_groups:
        delete_empty_group_shapes(slide, empty_groups)


def _fill_by_path(
    slide,
    prefix: Optional[str],
    content_map: Dict,
    type_map: Dict,
    used_shapes: Set[int],
    empty_groups: Set[str]
) -> None:
    """通过形状路径匹配填充内容。"""
    for shape, raw_path in iter_shapes_with_path(slide.shapes):
        label_path = normalize_path(raw_path, prefix)
        if not label_path:
            continue
        key = tuple(label_path)
        if key not in content_map:
            continue

        shape_name = shape.name or ""
        if any(kw in shape_name for kw in IGNORE_KEYWORDS):
            continue

        value = content_map[key]
        field_type = type_map.get(key)
        is_image = field_type == "image" if field_type else is_picture_shape(shape)

        if is_image:
            replace_picture(slide, shape, value)
            used_shapes.add(id(shape))
        elif shape.has_text_frame:
            if value:
                set_shape_text(shape, value)
            else:
                shape.text_frame.clear()
                if len(raw_path) > 1:
                    empty_groups.add("/".join(raw_path[:-1]))
            used_shapes.add(id(shape))


def _fill_by_key(
    slide,
    page_content: Dict,
    shapes_by_exact: Dict,
    shapes_by_name: Dict,
    shapes_by_tag: Dict,
    text_placeholders: List,
    picture_placeholders: List,
    used_shapes: Set[int],
    shape_name_to_group_path: Dict,
    empty_groups: Set[str],
    pop_placeholder_fn
) -> None:
    """通过键名匹配填充内容。"""
    for area_name, raw_value in page_content.items():
        field_type = None
        if isinstance(raw_value, dict):
            if "type" in raw_value and "value" in raw_value:
                value = raw_value.get("value", "")
                field_type = raw_value.get("type", "text")
            else:
                continue
        else:
            value = raw_value

        # 查找匹配的形状
        shape = _find_matching_shape(
            area_name, shapes_by_exact, shapes_by_name, shapes_by_tag,
            text_placeholders, picture_placeholders, used_shapes, pop_placeholder_fn
        )

        if not shape:
            print(f"⚠️  找不到名为「{area_name}」的形状，内容已忽略。")
            continue

        if id(shape) in used_shapes:
            continue

        is_image = field_type == "image" if field_type else is_picture_shape(shape)

        if is_image:
            replace_picture(slide, shape, value)
            used_shapes.add(id(shape))
            if not value and shape.name in shape_name_to_group_path:
                for gp in shape_name_to_group_path[shape.name]:
                    empty_groups.add(gp)
        elif shape.has_text_frame:
            if value:
                set_shape_text(shape, value)
            else:
                shape.text_frame.clear()
                if shape.name in shape_name_to_group_path:
                    for gp in shape_name_to_group_path[shape.name]:
                        empty_groups.add(gp)
            used_shapes.add(id(shape))
        else:
            print(f"⚠️  形状「{area_name}」既不是文本也不是图片，跳过。")


def _find_matching_shape(
    area_name: str,
    shapes_by_exact: Dict,
    shapes_by_name: Dict,
    shapes_by_tag: Dict,
    text_placeholders: List,
    picture_placeholders: List,
    used_shapes: Set[int],
    pop_placeholder_fn
):
    """查找与区域名称匹配的形状。"""
    # 精确匹配
    shape = take_shape(shapes_by_exact, area_name, used_shapes)

    # 别名匹配
    if not shape:
        for candidate in candidate_keys(area_name):
            shape = take_shape(shapes_by_name, candidate, used_shapes)
            if shape:
                break

    # 标签匹配
    if not shape:
        shape = take_shape(shapes_by_tag, area_name, used_shapes)
    if not shape:
        for candidate in candidate_keys(area_name):
            shape = take_shape(shapes_by_tag, candidate, used_shapes)
            if shape:
                break

    # 手动映射
    if not shape and area_name in MANUAL_NAME_MAP:
        for exact in MANUAL_NAME_MAP[area_name]:
            shape = take_shape(shapes_by_exact, exact, used_shapes)
            if shape:
                break
    if not shape and "字幕" in area_name:
        for exact in MANUAL_NAME_MAP.get("字幕", []):
            shape = take_shape(shapes_by_exact, exact, used_shapes)
            if shape:
                break

    # 占位符匹配
    if not shape and any(kw in area_name for kw in ("内容", "字幕")):
        shape = pop_placeholder_fn(text_placeholders, used_shapes)
    if not shape and any(kw in area_name for kw in ("图片区", "图片")):
        shape = pop_placeholder_fn(picture_placeholders, used_shapes)

    return shape


def build_from_json(
    template_path: Path,
    json_path: Path,
    output_path: Path
) -> None:
    """复制原始模板 pptx，并按 JSON 顺序重新组织 slide 文件。

    Args:
        template_path: 模板 PPTX 路径
        json_path: JSON 配置文件路径
        output_path: 输出 PPTX 路径
    """
    data = json.loads(Path(json_path).read_text(encoding="utf-8"))
    pages = data.get("ppt_pages", [])
    if not pages:
        raise ValueError("JSON 中未找到 ppt_pages 内容")

    with tempfile.TemporaryDirectory() as tmpdir:
        temp_copy = Path(tmpdir) / "working.pptx"
        shutil.copyfile(template_path, temp_copy)

        # 读取模板文件
        with zipfile.ZipFile(template_path, "r") as tmpl_zip:
            file_bytes = {name: tmpl_zip.read(name) for name in tmpl_zip.namelist()}

        slide_map, slide_rel_map = _extract_slides(file_bytes)
        slide_count = len(slide_map)
        if slide_count == 0:
            raise ValueError("模板中未找到任何 slide 文件")

        # 处理 tag 文件
        tag_nums = [
            int(m.group(1))
            for name in file_bytes
            if (m := TAG_RE.fullmatch(name)) is not None
        ]
        next_tag_num = max(tag_nums) if tag_nums else 0
        extra_tag_parts = []
        extra_tag_files = {}

        # 验证并打印处理进度
        selected_slides = _validate_pages(pages, slide_map)
        slide_total = len(selected_slides)

        # 更新 XML 结构
        pres_rels = ET.fromstring(file_bytes["ppt/_rels/presentation.xml.rels"])
        new_rel_ids = update_presentation_rels(pres_rels, slide_total)

        pres_xml = ET.fromstring(file_bytes["ppt/presentation.xml"])
        update_presentation_xml(pres_xml, new_rel_ids)

        # 处理 slide rels
        prepared_rel_bytes, next_tag_num, extra_tag_parts, extra_tag_files = (
            _prepare_slide_rels(
                selected_slides, slide_rel_map, file_bytes,
                next_tag_num, extra_tag_parts, extra_tag_files
            )
        )

        # 更新 Content Types
        content_types = ET.fromstring(file_bytes["[Content_Types].xml"])
        update_content_types(content_types, slide_total, extra_tag_parts)

        # 写入输出文件
        _write_output_pptx(
            output_path, file_bytes, slide_map, selected_slides,
            prepared_rel_bytes, extra_tag_files, content_types,
            pres_rels, pres_xml
        )

    print(f"\n🎉 新 PPT 输出完成：{output_path}")


def _extract_slides(file_bytes: Dict) -> Tuple[Dict, Dict]:
    """从文件字节中提取 slide 和 rels 映射。"""
    slide_map = {}
    slide_rel_map = {}
    for name in file_bytes:
        match = SLIDE_RE.fullmatch(name)
        if match:
            slide_map[int(match.group(1))] = file_bytes[name]
        elif name.startswith("ppt/slides/_rels/slide") and name.endswith(".xml.rels"):
            num = int(re.search(r"slide(\d+)\.xml\.rels", name).group(1))
            slide_rel_map[num] = file_bytes[name]
    return slide_map, slide_rel_map


def _validate_pages(pages: List[Dict], slide_map: Dict) -> List[Tuple[int, str]]:
    """验证 JSON 页面配置并返回选中的幻灯片列表。"""
    selected_slides = []
    for idx, page in enumerate(pages, start=1):
        tmpl_num = page.get("template_page_num")
        page_type = page.get("page_type", "未知版式")
        if tmpl_num is None:
            raise ValueError(f"第{idx}条缺少 template_page_num")
        if tmpl_num not in slide_map:
            raise ValueError(
                f"模板中不存在第{tmpl_num}页（来自第{idx}条 {page_type}）"
            )
        selected_slides.append((tmpl_num, page_type))
        print(f"✅ 生成第{idx}页：{page_type}（模板第{tmpl_num}页）")
    return selected_slides


def _prepare_slide_rels(
    selected_slides: List[Tuple[int, str]],
    slide_rel_map: Dict,
    file_bytes: Dict,
    next_tag_num: int,
    extra_tag_parts: List,
    extra_tag_files: Dict
) -> Tuple[Dict, int, List, Dict]:
    """准备 slide rels 文件。"""
    def clone_tags(rel_bytes: bytes) -> bytes:
        nonlocal next_tag_num
        if not rel_bytes:
            return rel_bytes
        rel_tree = ET.fromstring(rel_bytes)
        for rel in list(rel_tree.findall(f"{{{PKG_REL_NS}}}Relationship")):
            target = rel.get("Target", "")
            rel_type = rel.get("Type", "")

            if "notesSlide" in target or "notesSlide" in rel_type:
                rel_tree.remove(rel)
                continue

            if rel_type == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/tags":
                canonical = posixpath.normpath(posixpath.join("ppt/slides", target))
                if canonical not in file_bytes:
                    continue
                next_tag_num += 1
                new_part = f"ppt/tags/tag{next_tag_num}.xml"
                rel.set("Target", posixpath.relpath(new_part, "ppt/slides"))
                extra_tag_parts.append(new_part)
                extra_tag_files[new_part] = file_bytes[canonical]
        return clean_rels_namespace(
            ET.tostring(rel_tree, encoding="utf-8", xml_declaration=True)
        )

    prepared_rel_bytes = {}
    for idx, (tmpl_num, _) in enumerate(selected_slides, start=1):
        if tmpl_num in slide_rel_map:
            prepared_rel_bytes[idx] = clone_tags(slide_rel_map[tmpl_num])

    return prepared_rel_bytes, next_tag_num, extra_tag_parts, extra_tag_files


def _write_output_pptx(
    output_path: Path,
    file_bytes: Dict,
    slide_map: Dict,
    selected_slides: List[Tuple[int, str]],
    prepared_rel_bytes: Dict,
    extra_tag_files: Dict,
    content_types: ET.Element,
    pres_rels: ET.Element,
    pres_xml: ET.Element
) -> None:
    """写入输出 PPTX 文件。"""
    with zipfile.ZipFile(output_path, "w") as out_zip:
        # 写入非 slide 文件
        for name, data in file_bytes.items():
            if name.startswith("ppt/slides/slide"):
                continue
            if name.startswith("ppt/slides/_rels/slide"):
                continue
            if "notesSlides" in name or "notesMasters" in name:
                continue
            if name == "[Content_Types].xml":
                out_zip.writestr(
                    name,
                    ET.tostring(content_types, encoding="utf-8", xml_declaration=True),
                )
            elif name == "ppt/_rels/presentation.xml.rels":
                out_zip.writestr(
                    name,
                    clean_rels_namespace(
                        ET.tostring(pres_rels, encoding="utf-8", xml_declaration=True)
                    ),
                )
            elif name == "ppt/presentation.xml":
                out_zip.writestr(
                    name,
                    ET.tostring(pres_xml, encoding="utf-8", xml_declaration=True),
                )
            else:
                out_zip.writestr(name, data)

        # 写入 slide 文件
        for idx, (tmpl_num, _) in enumerate(selected_slides, start=1):
            slide_name = f"ppt/slides/slide{idx}.xml"
            rel_name = f"ppt/slides/_rels/slide{idx}.xml.rels"
            out_zip.writestr(slide_name, slide_map[tmpl_num])
            rel_bytes = prepared_rel_bytes.get(idx)
            if rel_bytes:
                out_zip.writestr(rel_name, rel_bytes)

        # 写入 tag 文件
        for name, data in extra_tag_files.items():
            out_zip.writestr(name, data)


def render_slides(
    template_path: Path,
    config: Dict,
    output_name: str,
    run_dir: Optional[Path] = None,
) -> Dict:
    """渲染入口，供 GUI/CLI 复用。

    Args:
        template_path: 模板 PPTX 路径
        config: JSON 配置字典
        output_name: 输出文件名
        run_dir: 输出目录（可选）

    Returns:
        包含 output_path, run_dir, slides 的字典
    """
    pages = config.get("ppt_pages", [])
    if not pages:
        raise ValueError("JSON 数据中没有 ppt_pages 内容。")

    run_dir = run_dir or create_run_dir()
    run_dir.mkdir(parents=True, exist_ok=True)
    output_path = run_dir / output_name

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        temp_ppt = Path(tmp.name)

    try:
        tmp_json = run_dir / "config.json"
        tmp_json.write_text(
            json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8"
        )

        build_from_json(template_path, tmp_json, temp_ppt)
        connector_snapshots = extract_connectors(temp_ppt)
        prs = Presentation(temp_ppt)
        if len(prs.slides) != len(pages):
            raise RuntimeError("生成的幻灯片数量与 JSON 不匹配，无法填充。")

        slide_width = prs.slide_width
        for slide, page in zip(prs.slides, pages):
            fill_slide(slide, page.get("content", {}), slide_width)

        prs.save(output_path)
        restore_connectors(output_path, connector_snapshots)
    finally:
        temp_ppt.unlink(missing_ok=True)

    return {"output_path": output_path, "run_dir": run_dir, "slides": len(pages)}
