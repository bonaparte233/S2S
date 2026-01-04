"""布局调整和自适应工具。"""

from pptx.util import Pt

from scripts.ppt_processing.constants import (
    SUBTITLE_TEXTS,
    PLACEHOLDER_KEYWORDS,
    EXPANDABLE_KEYWORDS,
    EMU_PER_PT,
    H_PADDING,
)
from scripts.ppt_processing.shape_utils import is_picture_shape


def clear_default_subtitles(shapes) -> None:
    """清除未被覆盖的字幕或默认说明文字。

    Args:
        shapes: 形状列表
    """
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if text in SUBTITLE_TEXTS or any(
            keyword in text for keyword in PLACEHOLDER_KEYWORDS
        ):
            shape.text_frame.clear()


def apply_layout_rules(shapes, slide_width: int) -> None:
    """根据文本长度自动调整标题条与背景。

    Args:
        shapes: 形状列表
        slide_width: 幻灯片宽度（EMU）
    """
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        name = shape.name or ""
        if not any(keyword in name for keyword in EXPANDABLE_KEYWORDS):
            continue
        _adjust_text_shape(shape, shapes, slide_width)


def _adjust_text_shape(text_shape, shapes, slide_width: int) -> None:
    """调整单个文本形状的宽度。"""
    text_width = _estimate_text_width(text_shape)
    if text_width <= 0:
        return

    limit = _find_right_limit(text_shape, shapes, slide_width) - H_PADDING
    available = max(text_shape.width, limit - text_shape.left)
    target_width = min(max(text_shape.width, text_width + H_PADDING), available)

    if target_width > text_shape.width:
        target_width = int(target_width)
        # 保存原始位置和高度
        original_left = text_shape.left
        original_top = text_shape.top
        original_height = text_shape.height
        # 修改宽度
        text_shape.width = target_width
        # 强制设置位置和高度，确保 xfrm 元素完整
        text_shape.left = original_left
        text_shape.top = original_top
        if original_height > 0:
            text_shape.height = original_height
        bg = _find_background_shape(text_shape, shapes)
        if bg:
            bg.width = int(max(bg.width, target_width + H_PADDING))
            bg.left = min(bg.left, text_shape.left)
    else:
        shrink_ratio = available / text_width if text_width else 1
        if shrink_ratio < 1:
            _shrink_font(text_shape, shrink_ratio)


def _find_background_shape(text_shape, shapes):
    """查找文本形状对应的背景形状。"""
    top = text_shape.top
    bottom = text_shape.top + text_shape.height
    candidate = None
    for shape in shapes:
        if shape is text_shape or shape.has_text_frame or is_picture_shape(shape):
            continue
        overlap = min(bottom, shape.top + shape.height) - max(top, shape.top)
        if overlap <= 0:
            continue
        ratio = overlap / max(1, text_shape.height)
        if ratio < 0.6:
            continue
        if candidate is None or shape.width > candidate.width:
            candidate = shape
    return candidate


def _find_right_limit(text_shape, shapes, slide_width: int) -> int:
    """查找文本形状右侧的边界限制。"""
    limit = slide_width
    top = text_shape.top
    bottom = text_shape.top + text_shape.height
    for shape in shapes:
        if shape is text_shape:
            continue
        other_top = shape.top
        other_bottom = shape.top + shape.height
        overlap = min(bottom, other_bottom) - max(top, other_top)
        if overlap <= 0:
            continue
        if shape.left > text_shape.left:
            limit = min(limit, shape.left)
    return limit


def _estimate_text_width(shape) -> int:
    """估算文本宽度（EMU）。"""
    if not shape.has_text_frame:
        return 0

    max_line = 0
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs) or para.text or ""
        if not line:
            continue
        font_size = None
        for run in para.runs:
            if run.font.size:
                font_size = run.font.size.pt
                break
        if font_size is None:
            font_size = 28
        width_factor = sum(0.55 if ord(ch) < 128 else 1 for ch in line)
        line_width = width_factor * font_size * EMU_PER_PT
        max_line = max(max_line, line_width)
    return max(max_line, shape.width)


def _shrink_font(shape, ratio: float) -> None:
    """缩小形状中的字体。"""
    ratio = max(ratio, 0.6)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                run.font.size = Pt(run.font.size.pt * ratio)
