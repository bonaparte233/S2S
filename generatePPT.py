"""根据 JSON 描述复制模板、填充文本/图片并生成最终 PPT。"""

import argparse
import json
import re
import tempfile
import zipfile
from pathlib import Path

from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from archive.generatePPT_template import build_from_json


SUFFIXES = ("区", "框", "栏")
PLACEHOLDER_KEYWORDS = ("文字内容", "字幕", "标题名称", "内容内容")
SUBTITLE_TEXTS = (
    "字幕18pt，白色字体深色描边，悬浮阴影。确保在任何底色上都能明确显示",
)
IGNORE_KEYWORDS = ("背景", "矩形", "圆角", "椭圆", "形状", "图形", "遮罩", "底色")
EXPANDABLE_KEYWORDS = ("标题", "名称", "课题", "栏目")
EMU_PER_PT = 12700
H_PADDING = 20000  # 约 1.5 毫米
MANUAL_NAME_MAP = {
    "目录内容区1": ["文本框 9"],
    "目录内容区2": ["文本框 14"],
    "目录内容区3": ["文本框 17"],
    "目录内容区4": ["文本框 20"],
    # 常见字幕框
    "字幕": ["文本框 10", "文本框 32", "文本框 36", "文本框 58", "文本框 121"],
}
NSMAP = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}
SLIDE_RE = re.compile(r"ppt/slides/slide(\d+)\.xml")


def _iter_shapes(shapes):
    """递归遍历幻灯片，包含组合内部的形状。"""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _iter_shapes_with_path(shapes, parent_path=None):
    """递归遍历幻灯片，返回 (shape, 路径)。"""
    for idx, shape in enumerate(shapes, start=1):
        name = (shape.name or "").strip() or f"元素{idx}"
        path = (*parent_path, name) if parent_path else (name,)
        yield shape, path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_with_path(shape.shapes, path)


def _is_picture_shape(shape):
    """判断形状是否充当图片占位符。"""
    name = shape.name or ""
    if "图片区" in name or name.startswith("图片"):
        return True
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if shape.is_placeholder:
        try:
            return shape.placeholder_format.type == 18
        except ValueError:
            return False
    return False


def _is_placeholder_shape(shape):
    """识别普通文本占位符，用于目录“文字文字”这类内容。"""
    if not shape.has_text_frame:
        return False
    text = (shape.text_frame.text or "").strip()
    name = shape.name or ""
    if name.startswith("文本框"):
        return True
    return any(keyword in text for keyword in PLACEHOLDER_KEYWORDS)


def _detect_prefix(slide):
    counter = {}
    for _, path in _iter_shapes_with_path(slide.shapes):
        for segment in path:
            prefix = _extract_prefix(segment)
            if prefix:
                counter[prefix] = counter.get(prefix, 0) + 1
    if not counter:
        return None
    return max(counter.items(), key=lambda item: item[1])[0]


def _extract_prefix(name):
    if "_" not in name:
        return None
    prefix = name.split("_", 1)[0]
    return prefix if "页" in prefix else None


def _normalize_path(path, page_prefix):
    start_idx = 0
    if page_prefix:
        start_idx = -1
        for idx, segment in enumerate(path):
            if page_prefix in segment:
                start_idx = idx
                break
        if start_idx == -1:
            return []
    trimmed = []
    for segment in path[start_idx:]:
        seg = _clean_segment(segment, page_prefix)
        if not seg:
            continue
        trimmed.append(seg)
    return trimmed


def _clean_segment(segment, page_prefix):
    seg = segment.strip()
    if not seg:
        return ""
    if page_prefix and seg.startswith(page_prefix + "_"):
        seg = seg[len(page_prefix) + 1 :]
    for keyword in IGNORE_KEYWORDS:
        if keyword in seg:
            return ""
    return seg


def _flatten_content(content):
    mapping = {}

    def walk(node, path):
        if isinstance(node, dict):
            for key, value in node.items():
                walk(value, path + (key,))
        else:
            mapping[path] = node

    walk(content or {}, tuple())
    return mapping


def _shape_aliases(name):
    """为形状名称生成多种别名，提升匹配成功率。"""
    aliases = set()
    clean = name.strip()
    if not clean:
        return aliases

    aliases.update({clean, clean.replace(" ", "")})

    def _add_parts(separator):
        if separator in clean:
            parts = [p for p in clean.split(separator) if p]
            for part in parts:
                aliases.add(part)
                aliases.add(part.replace(" ", ""))

    _add_parts("_")
    _add_parts("-")

    extra = set()
    for alias in aliases:
        for suf in SUFFIXES:
            if alias.endswith(suf):
                trimmed = alias[: -len(suf)]
                extra.add(trimmed)
                extra.add(trimmed.replace(" ", ""))
    aliases.update(extra)
    return aliases


def _candidate_keys(key):
    """给定 JSON 中的键名，生成若干匹配候选。"""
    key = key.strip()
    variants = [key, key.replace(" ", "")]
    for suf in SUFFIXES:
        if key.endswith(suf):
            variants.append(key[: -len(suf)])
            variants.append(key[: -len(suf)].replace(" ", ""))
    seen = set()
    result = []
    for variant in variants:
        if variant and variant not in seen:
            result.append(variant)
            seen.add(variant)
    return result


def _set_shape_text(shape, text):
    """填充文本时尽量保持原有格式。"""
    if not shape.has_text_frame:
        return

    text = "" if text is None else str(text)
    lines = text.split("\n")
    tf = shape.text_frame

    if not tf.paragraphs:
        tf.add_paragraph()

    for idx, line in enumerate(lines):
        if idx < len(tf.paragraphs):
            para = tf.paragraphs[idx]
        else:
            para = tf.add_paragraph()

        if para.runs:
            para.runs[0].text = line
            for run in para.runs[1:]:
                run.text = ""
        else:
            para.text = line

    # 清理多余段落
    for idx in range(len(lines), len(tf.paragraphs)):
        for run in tf.paragraphs[idx].runs:
            run.text = ""

    tf.word_wrap = True


def _safe_remove_shape(shape):
    element_parent = shape.element.getparent()
    if element_parent is not None:
        element_parent.remove(shape.element)


def _replace_picture(slide, shape, image_path):
    """将图片占位符替换为本地图片，保持位置大小比率。"""
    if not image_path:
        _safe_remove_shape(shape)
        return

    image_path = Path(image_path)
    if not image_path.is_file():
        print(f"⚠️  图片文件不可用：{image_path}")
        _safe_remove_shape(shape)
        return

    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    name = shape.name

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    img_ratio = img_w / img_h
    box_ratio = width / height
    if img_ratio > box_ratio:
        new_width = width
        new_height = width / img_ratio
    else:
        new_height = height
        new_width = height * img_ratio

    new_left = int(left + (width - new_width) / 2)
    new_top = int(top + (height - new_height) / 2)
    new_width = int(new_width)
    new_height = int(new_height)

    parent_shapes = getattr(shape, "_parent", None)
    _safe_remove_shape(shape)
    if parent_shapes is None:
        # 某些占位属于组合但 XML 已被剥离，此时无法安全移除，直接保留原父层
        parent_shapes = slide.shapes

    if parent_shapes is None or not hasattr(parent_shapes, "add_picture"):
        parent_shapes = slide.shapes

    new_pic = parent_shapes.add_picture(
        str(image_path), new_left, new_top, width=new_width, height=new_height
    )
    new_pic.name = name


def _fill_slide(slide, page_content, slide_width):
    """根据 JSON 内容把文本和图片写入对应区域。"""
    prefix = _detect_prefix(slide)
    content_map = _flatten_content(page_content)
    all_shapes = list(_iter_shapes(slide.shapes))
    shapes_by_name = {}
    shapes_by_exact = {}
    text_placeholders = []
    used_text_shapes = set()
    used_picture_shapes = set()
    picture_placeholders = []

    for shape in all_shapes:
        if shape.name:
            shapes_by_exact.setdefault(shape.name, shape)
            for alias in _shape_aliases(shape.name):
                shapes_by_name.setdefault(alias, shape)
            if _is_placeholder_shape(shape):
                text_placeholders.append(shape)
            if _is_picture_shape(shape):
                picture_placeholders.append(shape)

    for shape, raw_path in _iter_shapes_with_path(slide.shapes):
        label_path = _normalize_path(raw_path, prefix)
        if not label_path:
            continue
        key = tuple(label_path)
        if key not in content_map:
            continue
        value = content_map[key]

        if _is_picture_shape(shape):
            _replace_picture(slide, shape, value)
            used_picture_shapes.add(shape.name)
        elif shape.has_text_frame:
            if value:
                _set_shape_text(shape, value)
                used_text_shapes.add(shape.name)
            else:
                shape.text_frame.clear()
                used_text_shapes.add(shape.name)
        else:
            continue

    # 兼容旧 JSON 的键名，用别名机制兜底
    for area_name, value in page_content.items():
        if isinstance(value, dict):
            continue
        shape = None
        for candidate in _candidate_keys(area_name):
            shape = shapes_by_name.get(candidate)
            if shape:
                break
        if not shape and area_name in MANUAL_NAME_MAP:
            for exact in MANUAL_NAME_MAP[area_name]:
                shape = shapes_by_exact.get(exact)
                if shape:
                    break
        if not shape and "字幕" in area_name:
            for exact in MANUAL_NAME_MAP.get("字幕", []):
                shape = shapes_by_exact.get(exact)
                if shape:
                    break

        if not shape and any(keyword in area_name for keyword in ("内容", "字幕")):
            if text_placeholders:
                shape = text_placeholders.pop(0)
        if not shape and any(keyword in area_name for keyword in ("图片区", "图片")):
            if picture_placeholders:
                shape = picture_placeholders.pop(0)

        if not shape:
            print(f"⚠️  找不到名为“{area_name}”的形状，内容已忽略。")
            continue

        if _is_picture_shape(shape):
            _replace_picture(slide, shape, value)
        elif shape.has_text_frame:
            if value:
                _set_shape_text(shape, value)
            else:
                shape.text_frame.clear()
        else:
            print(f"⚠️  形状“{area_name}”既不是文本也不是图片，跳过。")

    _clear_default_subtitles(all_shapes)
    _apply_layout_rules(all_shapes, slide_width)


def _clear_default_subtitles(shapes):
    """清除未被覆盖的字幕或默认说明文字。"""
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if text in SUBTITLE_TEXTS or any(keyword in text for keyword in PLACEHOLDER_KEYWORDS):
            shape.text_frame.clear()


def _apply_layout_rules(shapes, slide_width):
    """根据文本长度自动调整标题条与背景。"""
    for shape in shapes:
        if not shape.has_text_frame:
            continue
        name = shape.name or ""
        if not any(keyword in name for keyword in EXPANDABLE_KEYWORDS):
            continue
        _adjust_text_shape(shape, shapes, slide_width)


def _adjust_text_shape(text_shape, shapes, slide_width):
    text_width = _estimate_text_width(text_shape)
    if text_width <= 0:
        return

    limit = _find_right_limit(text_shape, shapes, slide_width) - H_PADDING
    available = max(text_shape.width, limit - text_shape.left)
    target_width = min(max(text_shape.width, text_width + H_PADDING), available)

    if target_width > text_shape.width:
        target_width = int(target_width)
        delta = target_width - text_shape.width
        text_shape.width = target_width
        bg = _find_background_shape(text_shape, shapes)
        if bg:
            bg.width = int(max(bg.width, target_width + H_PADDING))
            bg.left = min(bg.left, text_shape.left)
    else:
        shrink_ratio = available / text_width if text_width else 1
        if shrink_ratio < 1:
            _shrink_font(text_shape, shrink_ratio)


def _find_background_shape(text_shape, shapes):
    top = text_shape.top
    bottom = text_shape.top + text_shape.height
    candidate = None
    for shape in shapes:
        if shape is text_shape or shape.has_text_frame or _is_picture_shape(shape):
            continue
        overlap = min(bottom, shape.top + shape.height) - max(top, shape.top)
        if overlap <= 0:
            continue
        ratio = overlap / max(1, text_shape.height)
        if ratio < 0.6:
            continue
        if candidate is None or shape.width > candidate.width:
            candidate = shape
    return candidate


def _find_right_limit(text_shape, shapes, slide_width):
    limit = slide_width
    top = text_shape.top
    bottom = text_shape.top + text_shape.height
    for shape in shapes:
        if shape is text_shape:
            continue
        other_top = shape.top
        other_bottom = shape.top + shape.height
        overlap = min(bottom, other_bottom) - max(top, other_top)
        if overlap <= 0:
            continue
        if shape.left > text_shape.left:
            limit = min(limit, shape.left)
    return limit


def _estimate_text_width(shape):
    if not shape.has_text_frame:
        return 0

    max_line = 0
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text for run in para.runs) or para.text or ""
        if not line:
            continue
        font_size = None
        for run in para.runs:
            if run.font.size:
                font_size = run.font.size.pt
                break
        if font_size is None:
            font_size = 28
        width_factor = sum(0.55 if ord(ch) < 128 else 1 for ch in line)
        line_width = width_factor * font_size * EMU_PER_PT
        max_line = max(max_line, line_width)
    return max(max_line, shape.width)


def _shrink_font(shape, ratio):
    ratio = max(ratio, 0.6)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size:
                run.font.size = Pt(run.font.size.pt * ratio)


def _extract_connectors(pptx_path: Path):
    connectors = {}
    with zipfile.ZipFile(pptx_path, "r") as zf:
        for name in zf.namelist():
            match = SLIDE_RE.fullmatch(name)
            if not match:
                continue
            slide_idx = int(match.group(1))
            root = etree.fromstring(zf.read(name))
            nodes = root.xpath(".//p:cxnSp", namespaces=NSMAP)
            if nodes:
                connectors[slide_idx] = [etree.tostring(node) for node in nodes]
    return connectors


def _restore_connectors(pptx_path: Path, connectors):
    if not connectors:
        return
    with zipfile.ZipFile(pptx_path, "r") as src:
        entries = {name: src.read(name) for name in src.namelist()}
    modified = False
    for name, data in list(entries.items()):
        match = SLIDE_RE.fullmatch(name)
        if not match:
            continue
        slide_idx = int(match.group(1))
        snippets = connectors.get(slide_idx)
        if not snippets:
            continue
        root = etree.fromstring(data)
        sp_tree = root.find(".//p:spTree", namespaces=NSMAP)
        if sp_tree is None:
            continue
        for node in sp_tree.findall("p:cxnSp", namespaces=NSMAP):
            sp_tree.remove(node)
        for snippet in snippets:
            sp_tree.append(etree.fromstring(snippet))
        entries[name] = etree.tostring(root, encoding="utf-8", xml_declaration=True)
        modified = True
    if not modified:
        return
    with zipfile.ZipFile(pptx_path, "w") as dst:
        for name, data in entries.items():
            dst.writestr(name, data)


def generate_ppt(template_path, json_path, output_path):
    """完整流程：复制模板顺序 -> 填充内容 -> 输出 PPT。"""
    pages = json.loads(Path(json_path).read_text(encoding="utf-8")).get("ppt_pages", [])
    if not pages:
        raise ValueError("JSON 文件中没有 ppt_pages 内容。")

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        temp_ppt = Path(tmp.name)

    try:
        build_from_json(template_path, json_path, temp_ppt)
        connector_snapshots = _extract_connectors(temp_ppt)
        prs = Presentation(temp_ppt)
        if len(prs.slides) != len(pages):
            raise RuntimeError("生成的幻灯片数量与 JSON 不匹配，无法填充。")

        slide_width = prs.slide_width
        for slide, page in zip(prs.slides, pages):
            _fill_slide(slide, page.get("content", {}), slide_width)

        prs.save(output_path)
        _restore_connectors(Path(output_path), connector_snapshots)
    finally:
        temp_ppt.unlink(missing_ok=True)

    print(f"\n🎯 已根据内容生成 PPT：{output_path}")


def main():
    parser = argparse.ArgumentParser(description="读取 JSON 并填充模板生成 PPT")
    parser.add_argument("--template", required=True, help="模板 PPTX 路径")
    parser.add_argument("--json", required=True, help="描述内容的 JSON 文件")
    parser.add_argument("--output", default="final_output.pptx", help="输出 PPTX 路径")
    args = parser.parse_args()

    generate_ppt(args.template, args.json, args.output)


if __name__ == "__main__":
    main()
