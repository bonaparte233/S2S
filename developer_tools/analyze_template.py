"""
开发者工具: 半自动生成 map.json 草稿

功能: 分析 PowerPoint 模板，提取布局和占位符信息，生成 map.json 草稿文件
注意: 生成的草稿文件中的 constraints 需要开发者手动填写具体数值
"""

import argparse
import json
import os
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


def generate_map_draft(template_path, output_path):
    """
    生成 map.json 草稿文件

    :param template_path: PowerPoint 模板文件路径
    :param output_path: 输出的 map.json 文件路径
    """
    # 输入验证
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"模板文件不存在: {template_path}")

    # 加载 PowerPoint 模板
    prs = Presentation(template_path)

    # 初始化模板映射字典
    template_map = {"template_file": os.path.basename(template_path), "layouts": {}}

    # 全局布局索引计数器
    global_layout_index = 0

    # 检查模板是否有示例幻灯片
    has_template_slides = len(prs.slides) > 0

    # 遍历所有 Slide Masters（一个模板可能有多个 Master）
    for master_idx, slide_master in enumerate(prs.slide_masters):
        # 遍历当前 Master 下的所有布局
        for local_layout_index, layout in enumerate(slide_master.slide_layouts):
            layout_index = global_layout_index
            global_layout_index += 1

            # 使用布局索引作为 map.json 的键
            layout_id_str = str(layout_index)

            # 初始化布局条目
            layout_entry = {
                "name": layout.name,
                "layout_index": layout_index,
                "master_index": master_idx,
                "placeholders": {},
            }

            # 如果模板有示例幻灯片，且索引在范围内，记录对应的幻灯片索引
            if has_template_slides and layout_index < len(prs.slides):
                layout_entry["template_slide_index"] = layout_index

            # 计数器，用于处理同类型的多个占位符
            body_count = 0
            picture_count = 0
            unknown_count = 0

            # 遍历布局中的所有占位符
            for ph in layout.placeholders:
                ph_format = ph.placeholder_format
                ph_id = ph_format.idx
                ph_type = ph_format.type

                # 跳过不需要的占位符类型（日期、页脚、页码等）
                skip_types = [
                    PP_PLACEHOLDER.DATE,
                    PP_PLACEHOLDER.FOOTER,
                    PP_PLACEHOLDER.SLIDE_NUMBER,
                ]
                if ph_type in skip_types:
                    continue

                # 启发式规则: 根据占位符类型和名称推断键名
                if ph_type == PP_PLACEHOLDER.TITLE:
                    key_name = "title"
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    key_name = "subtitle"
                elif ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                    body_count += 1
                    key_name = (
                        "body_text" if body_count == 1 else f"body_text_{body_count}"
                    )
                elif ph_type == PP_PLACEHOLDER.PICTURE:
                    picture_count += 1
                    key_name = (
                        "picture" if picture_count == 1 else f"picture_{picture_count}"
                    )
                else:
                    # 对于未知类型，尝试从名称推断
                    ph_name_lower = ph.name.lower()
                    if "title" in ph_name_lower or "标题" in ph_name_lower:
                        key_name = "title"
                    elif "subtitle" in ph_name_lower or "副标题" in ph_name_lower:
                        key_name = "subtitle"
                    elif (
                        "picture" in ph_name_lower
                        or "图片" in ph_name_lower
                        or "图像" in ph_name_lower
                    ):
                        picture_count += 1
                        key_name = (
                            "picture"
                            if picture_count == 1
                            else f"picture_{picture_count}"
                        )
                    elif (
                        "text" in ph_name_lower
                        or "内容" in ph_name_lower
                        or "正文" in ph_name_lower
                    ):
                        body_count += 1
                        key_name = (
                            "body_text"
                            if body_count == 1
                            else f"body_text_{body_count}"
                        )
                    else:
                        unknown_count += 1
                        key_name = f"unknown_{ph_id}_({ph.name})"

                # 确保键名唯一性
                original_key_name = key_name
                suffix = 2
                while key_name in layout_entry["placeholders"]:
                    key_name = f"{original_key_name}_{suffix}"
                    suffix += 1

                # 生成占位符条目（包含详细信息以提高可读性）
                ph_entry = {
                    "id": ph_id,
                    "type": str(ph_format.type),  # 占位符类型（便于理解）
                    "name": ph.name,  # 占位符名称
                    "left": ph.left,  # 位置信息（便于识别）
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height,
                    "constraints": {
                        "max_chars": "50",
                        "max_lines": "2",
                    },
                }

                layout_entry["placeholders"][key_name] = ph_entry

            # 将布局条目添加到模板映射
            template_map["layouts"][layout_id_str] = layout_entry

    # 确保输出目录存在
    output_dir = os.path.dirname(output_path)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)

    # 输出 JSON 文件
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(template_map, f, indent=2, ensure_ascii=False)

    print(f"✓ 成功生成 map.json 草稿: {output_path}")
    print("\n⚠️  重要提示:")
    print("   请**务必手动编辑**此文件，完成以下操作:")
    print("   1. 审查所有占位符的键名 (key_name)，确保其语义清晰")
    print("   2. 用具体的**数字**替换所有 'AUTO_GENERATED_PLEASE_FILL'")
    print("   3. 根据实际模板设置合理的 max_chars 和 max_lines 约束")


if __name__ == "__main__":
    # 命令行参数解析
    parser = argparse.ArgumentParser(
        description="分析 PowerPoint 模板并生成 map.json 草稿文件"
    )
    parser.add_argument(
        "-t",
        "--template",
        required=True,
        help="PowerPoint 模板文件路径 (例如: inputs/templates/template.pptx)",
    )
    parser.add_argument(
        "-o",
        "--output",
        required=True,
        help="输出的 map.json 文件路径 (例如: config_maps/template.map.json)",
    )

    args = parser.parse_args()

    # 执行生成
    try:
        generate_map_draft(args.template, args.output)
    except Exception as e:
        import traceback

        print(f"✗ 错误: {e}")
        print("\n详细错误信息:")
        traceback.print_exc()
        exit(1)
